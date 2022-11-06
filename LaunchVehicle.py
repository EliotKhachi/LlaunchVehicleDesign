# NOTES:
# Must still modify avionics and wiring estimates for small LVs

from Mission import Mission

from math import sqrt, pow, pi, floor, ceil
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import ColorFormat, RGBColor
from pptx.util import Inches
import matlab.engine

class LaunchVehicle(Mission) :
    
    def __init__(self, name, TW, body_material, num_steps, engine_Isps, sigmas, PL, Mission):
        self.name = name
        self.TW = TW
        self.body_material = body_material
        if body_material == "Aluminum 6061-T6":
            self.rho_body = 2700  # kg/m^3 al 6061-T6  density
        elif body_material == "Aluminum 7075-T6":
            self.rho_body = 2810 # kg/m^3 al 7075-T6 density
        elif body_material == "Aluminum 2219-T87":
            self.rho_body = 2840 # kg/m^3 al 7075-T6 density
        self.num_steps = num_steps
        self.engine_Isps = engine_Isps
        self.sigmas = sigmas
        self.PL = PL
        self.Mission = Mission
        print(self.name + ' has been initialized')
        if self.Mission.input[0] == 'One':
            self.m_PAF = 7 # CAD estimate of PAF mass
        elif self.Mission.input[0] == 'Two':
            self.m_PAF = 14 # CAD estimate of PAF mass
        self.m_engine_0 = 22 # Edberg's textbook
        # LR101 Engine Info (22kg weight) http://astronautix.com/l/lr101-11.html


    def initMassEstimates(self):
        df_masses = pd.read_csv('LVMasses\\' + self.name + 'MassEstimate.csv', index_col=0)
        # Initialize m_gross of the Launch Vehicle
        self.m_gross = df_masses.iloc[6, 2]

        # Initialize propellant masses (list of masses per step)
        self.m_p = list(df_masses.iloc[3, :].astype(float)) # get list from panda series
        self.m_p.reverse() # reverse list order (from steps bottom to top)
        self.m_p = [item for item in self.m_p if item != 0] # remove 0 values

        # Initialize structural masses (list of masses per step)
        self.m_s = list(df_masses.iloc[4, :].astype(float)) # get list from panda series
        self.m_s.reverse() # reverse list order (from steps bottom to top)
        self.m_s = [item for item in self.m_s if item != 0] # remove 0 values

        # Initialize stage masses (list of masses per step)
        self.m_0 = list(df_masses.iloc[5, :].astype(float)) # get list from panda series
        self.m_0.reverse() # reverse list order (from steps bottom to top)
        self.m_0 = [item for item in self.m_0 if item != 0] # remove 0 values

        print(self.m_0)
        print(self.m_s)
        print(self.m_p)

        self.mp_actual = [] # initialized in initMasses
        self.mi_actual = [] # initialized in rearrangeDF
        self.mf_actual = [] # initialized in rearrangeDF
        print("the gross mass is " + str (self.m_gross))
        print("the stage masses are " + str(self.m_0))
        print("the propellant masses are " + str(self.m_p))
        print("the structural masses are " + str(self.m_s))

    def initSteps(self, listOfSteps):
        print("Initializing and sizing steps...")
        self.listOfSteps = listOfSteps
        self.r = []
        for i in range(len(self.listOfSteps)):
            step = listOfSteps[i]
            self.r.append(step.r)

    def initInterstages(self): #  adds interstage to Step object as its "Forward Skirt"
        pi = 3.1415926535897932#3846264338327950288
        #self.listOfInterstages = []
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            temp_interstage = []
            if i == len(self.listOfSteps)-1: # if last step: last step doesn't have interstage
                temp_interstage.append(0)
                temp_interstage.append(0)
            else: # if not last step
                if step.propulsion == 'Liquid':
                    temp_interstage.append(step.dome_f[0] + self.listOfSteps[i+1].L_n + step.r/4) # interstage length: gap of r/4 between the current step's tank dome and upper step's nozzle exit
                    temp_interstage.append(pi*(step.r + self.listOfSteps[i+1].r)* pow( pow(step.r - self.listOfSteps[i+1].r, 2) + pow(temp_interstage[0], 2), 1/2) ) # interstage surface area -> thin frustum of cone
                elif step.propulsion == 'Solid':
                    if step.parallel:
                        temp_interstage.append(2 * step.press_tank[0] + step.r/2) # "interstage" length is actually the fwd skirt of the booster: total height of pressure tank + gap of r/2
                        temp_interstage.append(2 * pi * step.r * temp_interstage[0]) # interstage surface area -> thin cylinder
                    else:
                        temp_interstage.append(2 * step.press_tank[0] + self.listOfSteps[i+1].L_n + step.r/4) # interstage length: current step's total pressure tank length + gap of r/4 + next step's nozzle length
                        temp_interstage.append(pi*(step.r + self.listOfSteps[i+1].r)* pow( pow(step.r - self.listOfSteps[i+1].r, 2) + pow(temp_interstage[0], 2), 1/2) ) # interstage surface area -> thin frustum of cone
                        
            step.interstage = temp_interstage # set temp_interstage as the step's 'interstage' attribute

        #self.listOfSteps[0].eng_protrude_dist = self.listOfSteps[0].L_n/2
        
    def massMoments(self, load_type):
        print("Generating Mass Moments Table...")
        payload_items = ['TopOfPLF', 'PLF', 'Payload', 'PAF']
        step_items_liquid = ['Forward Skirt', 'Avionics', 'Wiring', 'Fuel Dome Top', 'Fuel Cylinder', 'Fuel Dome Bottom', 'Fuel Insulation', 'Fuel Residual', 
                             'Intertank', 'Ox Dome Top', 'Ox Cylinder', 'Ox Dome Bottom', 'Ox Insulation', 'Ox Residual', 'Pressurant Tank', 'Aft Skirt', 'Thrust Structure', 
                             'Gimballs', 'Engines', 'Fuel', 'Oxidizer']

        step_items_solid =['Nose Cone', 'Forward Skirt', 'Avionics', 'Wiring', 'Pressurant Tank', 'SRM Dome Top', 'Solid Propellant Casing', 'SRM Dome Bottom', 'Solid Propellant Residual', 'Aft Skirt',
                           'Gimballs', 'Nozzle', 'Solid Propellant']
        
        #num_items = len(payload_items) + self.num_steps * len(step_items)
        df_temp = pd.DataFrame(columns=['Item', 'Height (m)', 'Mass (kg)', 'Distance (m)', 'Moment (kg*m)', 'Thickness (m)', 'Distance from CM (m)', 'J0 (kg m^2)', 'm*CM^2 (kg m^2)', 'Jpitch/yaw', 'Jroll'],
                          index=range(len(payload_items)))
        #df = pd.DataFrame(columns=['Item', 'Height (m)', 'Mass (kg)', 'Distance (m)', 'Moment (kg*m)', 'Thickness (m)', 'Distance from CM (m)', 'J0 (kg m^2)', 'm*CM^2 (kg m^2)', 'Jpitch/yaw', 'Jroll'],
        #                  index=range(num_items))

        self.df = self.appendItems(df_temp, payload_items, step_items_liquid, step_items_solid)
        self.initHeights(self.df, payload_items, step_items_liquid, step_items_solid)
        self.initThicknesses(self.df, payload_items, step_items_liquid, step_items_solid)
        self.initMasses(self.df, payload_items, step_items_liquid, step_items_solid, load_type)
        self.initDistances(self.df, payload_items, step_items_liquid, step_items_solid, load_type)
        self.initMoments(self.df, payload_items, step_items_liquid, step_items_solid)
        self.initJ0s(self.df, payload_items, step_items_liquid, step_items_solid)
        #print("The sum of Moments is " + str(self.df['Moment (kg*m)'].sum()))
        #print("The sum of Masses is " + str(self.df['Mass (kg)'].sum()))
        self.CM_full = self.df['Moment (kg*m)'].sum()/self.df['Mass (kg)'].sum()
        print("The CM fully loaded of " + self.name + " is " + str(self.CM_full))
        self.initDistFromCM(self.df, payload_items, step_items_liquid, step_items_solid)
        self.initmCMs(self.df, payload_items, step_items_liquid, step_items_solid)
        self.initJPitchYaw(self.df, payload_items, step_items_liquid, step_items_solid)
        self.initJRoll(self.df, payload_items, step_items_liquid, step_items_solid)
        self.rearrangeDF(payload_items, step_items_liquid, step_items_solid) # must be called after initMassesLV and initMasses
        if load_type == 'Ground Wind-Loads Condition':
            self.df.to_csv('LVMassMoments\\'+self.name + 'GroundLoadsMassMoments.csv')
        elif load_type == 'Max-Q Condition':
            self.df.to_csv('LVMassMoments\\'+self.name + 'WindLoadsMassMoments.csv')
        for i in self.listOfSteps:
            
            print("step " + str(i.step_num) + " has a radius of " + str(i.r) + " m.")
            #print("thrust sea_level is " + str(i.T_SL) + " and step mass is " + str(sum(self.m_0)))

    def appendItems(self, df, pl_items, l_items, s_items):# initialize 'Item' names
        for i in range(len(pl_items)): # initialize payload item names
            df['Item'][i] = pl_items[i]
        for i in range(len(self.listOfSteps)): # initialize step item names
            step = self.listOfSteps[i]
            if step.propulsion == 'Liquid': # if liquid step
                for j in range(len(l_items)):
                    new_row = {'Item': l_items[j] + ' ' + str(step.step_num)}
                    df = df.append(new_row, ignore_index=True)
                    
            elif step.propulsion == 'Solid': # if solid step
                for j in range(len(s_items)):
                    new_row = {'Item': s_items[j] + ' ' + str(step.step_num)}
                    df = df.append(new_row, ignore_index=True)
        return df
    
    def initHeights(self, df, pl_items, l_items, s_items ):
        #print(self.listOfSteps)
        df['Height (m)'][pl_items.index("PLF")] = sum(self.listOfSteps[len(self.listOfSteps)-1].fairing[0])
        df['Height (m)'][pl_items.index("Payload")] = sum(self.listOfSteps[len(self.listOfSteps)-1].fairing[0]) 
        df['Height (m)'][pl_items.index("PAF")] = " "
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            
            if step.propulsion == 'Liquid':

                if i < len(self.listOfSteps) - 1: # if not the last step
                    df['Height (m)'][row + l_items.index("Forward Skirt")] = step.interstage[0] # Interstage 1 (still called fwd skirt)
                    #print("The forward skirt length is " + str(step.interstage[0]))
                
                elif i == len(self.listOfSteps) - 1: #else if last step
                    df['Height (m)'][row + l_items.index("Forward Skirt")] = step.fwd_skirt[0] # fwd skirt
                    #print("The forward skirt length is " + str(step.fwd_skirt[0]))
                #print("The fuel cyl length is " + str(step.cyl_f[0]))
                df['Height (m)'][row + l_items.index("Avionics")] = " "                
                df['Height (m)'][row + l_items.index("Wiring")] = step.total_length
                df['Height (m)'][row + l_items.index("Fuel Dome Top")] = step.dome_f[0]      
                df['Height (m)'][row + l_items.index("Fuel Cylinder")] = step.cyl_f[0]        
                df['Height (m)'][row + l_items.index("Fuel Dome Bottom")] = step.dome_f[0]        
                df['Height (m)'][row + l_items.index("Fuel Insulation")] = step.cyl_f[0] + 2 * step.dome_f[0]   
                df['Height (m)'][row + l_items.index("Fuel Residual")] = " "              
                df['Height (m)'][row + l_items.index("Intertank")] = step.intertank[0]            
                df['Height (m)'][row + l_items.index("Ox Dome Top")] = step.dome_ox[0]    
                df['Height (m)'][row + l_items.index("Ox Cylinder")] = step.cyl_ox[0]  
                df['Height (m)'][row + l_items.index("Ox Dome Bottom")] = step.dome_ox[0]
                df['Height (m)'][row + l_items.index("Ox Insulation")] = step.cyl_ox[0] + 2 * step.dome_ox[0]
                df['Height (m)'][row + l_items.index("Ox Residual")] = " "                            
                df['Height (m)'][row + l_items.index("Pressurant Tank")] = 2 * step.press_tank[0]                
                df['Height (m)'][row + l_items.index("Aft Skirt")] = step.aft_skirt[0]                 
                df['Height (m)'][row + l_items.index("Thrust Structure")] = step.T_struct            
                df['Height (m)'][row + l_items.index("Gimballs")] = " "               
                df['Height (m)'][row + l_items.index("Engines")] = step.L_n      
                df['Height (m)'][row + l_items.index("Fuel")] = step.cyl_f[0] + 2 * step.dome_f[0] 
                df['Height (m)'][row + l_items.index("Oxidizer")] = step.cyl_ox[0] + 2 * step.dome_ox[0] 
                row += len(l_items)
                
            elif step.propulsion == 'Solid': 
                if i < len(self.listOfSteps) - 1: # if not the last step
                    if step.parallel:
                        df['Height (m)'][row + s_items.index("Nose Cone")] = step.fairing[0]
                    df['Height (m)'][row + s_items.index("Forward Skirt")] = step.interstage[0] 
                
                elif i == len(self.listOfSteps) - 1: #else if last step
                    df['Height (m)'][row + s_items.index("Forward Skirt")] = step.fwd_skirt[0] 

                df['Height (m)'][row + s_items.index("Avionics")] = " "    
                df['Height (m)'][row + s_items.index("Wiring")] = step.total_length
                df['Height (m)'][row + s_items.index("Pressurant Tank")] = 2 * step.press_tank[0]
                #print(step.dome_f[0])
                df['Height (m)'][row + s_items.index("SRM Dome Top")] = step.dome_f[0]
                df['Height (m)'][row + s_items.index("Solid Propellant Casing")] = step.srm_casing[0]
                df['Height (m)'][row + s_items.index("SRM Dome Bottom")] = step.dome_f[0]
                df['Height (m)'][row + s_items.index("Solid Propellant Residual")] = " "  
                df['Height (m)'][row + s_items.index("Aft Skirt")] = step.aft_skirt[0]
                df['Height (m)'][row + s_items.index("Gimballs")] = " "
                df['Height (m)'][row + s_items.index("Nozzle")] = step.L_n
                df['Height (m)'][row + s_items.index("Solid Propellant")] = step.srm_casing[0]
                row += len(s_items)
    
    def initThicknesses(self, df, pl_items, l_items, s_items):
        df['Thickness (m)'][pl_items.index("PLF")] = 1
        df['Thickness (m)'][pl_items.index("Payload")] = 0  
        df['Thickness (m)'][pl_items.index("PAF")] = 0  
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            if step.propulsion == 'Liquid':
                if i < len(self.listOfSteps) - 1:
                    df['Thickness (m)'][row + l_items.index("Forward Skirt")] = 1
                
                elif i == len(self.listOfSteps) - 1:
                    df['Thickness (m)'][row + l_items.index("Forward Skirt")] = 1

                df['Thickness (m)'][row + l_items.index("Avionics")] = 0
                df['Thickness (m)'][row + l_items.index("Wiring")] = 0
                df['Thickness (m)'][row + l_items.index("Pressurant Tank")] = 1
                df['Thickness (m)'][row + l_items.index("Fuel Dome Top")] = 1
                df['Thickness (m)'][row + l_items.index("Fuel Cylinder")] = 1
                df['Thickness (m)'][row + l_items.index("Fuel Dome Bottom")] = 1 
                df['Thickness (m)'][row + l_items.index("Fuel Insulation")] = 0
                df['Thickness (m)'][row + l_items.index("Fuel Residual")] = 0
                df['Thickness (m)'][row + l_items.index("Intertank")] = 1
                df['Thickness (m)'][row + l_items.index("Ox Dome Top")] = 1
                df['Thickness (m)'][row + l_items.index("Ox Cylinder")] = 1 
                df['Thickness (m)'][row + l_items.index("Ox Dome Bottom")] = 1
                df['Thickness (m)'][row + l_items.index("Ox Insulation")] = 0
                df['Thickness (m)'][row + l_items.index("Ox Residual")] = 0
                df['Thickness (m)'][row + l_items.index("Pressurant Tank")] = 1
                df['Thickness (m)'][row + l_items.index("Aft Skirt")] = 1
                df['Thickness (m)'][row + l_items.index("Thrust Structure")] = 0
                df['Thickness (m)'][row + l_items.index("Gimballs")] = 0
                df['Thickness (m)'][row + l_items.index("Engines")] = 0 
                df['Thickness (m)'][row + l_items.index("Fuel")] = 0
                df['Thickness (m)'][row + l_items.index("Oxidizer")] = 0
                if ((self.name == 'Minerva-1') | (self.name == 'Minerva-2')) & (step.step_num == 1): # POST BUCKLING ANALYSIS CHANGES
                    df['Thickness (m)'][row + l_items.index("Ox Cylinder")] = 1.5
                    df['Thickness (m)'][row + l_items.index("Fuel Cylinder")] = 1.5
                    df['Thickness (m)'][row + l_items.index("Aft Skirt")] = 1.9
                elif ((self.name == 'Latona-1') & (step.step_num == 1)):
                    df['Thickness (m)'][row + l_items.index("Aft Skirt")] = 1.5
                elif ((self.name == 'Latona-2') & (step.step_num == 2)):
                    df['Thickness (m)'][row + l_items.index("Aft Skirt")] = 1.5
                elif ((self.name == 'Zephyr-1') | (self.name == 'Zephyr-2')) & (step.step_num == 1):
                    df['Thickness (m)'][row + l_items.index('Forward Skirt')] = 1.19
                    df['Thickness (m)'][row + l_items.index('Fuel Cylinder')] = 1.08
                    df['Thickness (m)'][row + l_items.index('Ox Cylinder')] = 1.42
                    df['Thickness (m)'][row + l_items.index('Intertank')] = 1.51
                    df['Thickness (m)'][row + l_items.index('Aft Skirt')] = 1.95
                row += len(l_items)
            elif step.propulsion == 'Solid':
                if i < len(self.listOfSteps) - 1:
                    if step.parallel:
                        df['Thickness (m)'][row + s_items.index("Nose Cone")] = 1
                    df['Thickness (m)'][row + s_items.index("Forward Skirt")] = 1
                
                elif i == len(self.listOfSteps) - 1:
                    df['Thickness (m)'][row + s_items.index("Forward Skirt")] = 1

                df['Thickness (m)'][row + s_items.index("Avionics")] = 0
                df['Thickness (m)'][row + s_items.index("Wiring")] = 0
                df['Thickness (m)'][row + s_items.index("Pressurant Tank")] = 1
                df['Thickness (m)'][row + s_items.index("SRM Dome Top")] = 1
                df['Thickness (m)'][row + s_items.index("Solid Propellant Casing")] = 1
                df['Thickness (m)'][row + s_items.index("SRM Dome Bottom")] = 1
                df['Thickness (m)'][row + s_items.index("Solid Propellant Residual")] = 0
                df['Thickness (m)'][row + s_items.index("Aft Skirt")] = 1
                df['Thickness (m)'][row + s_items.index("Gimballs")] = 0
                df['Thickness (m)'][row + s_items.index("Nozzle")] = 0
                df['Thickness (m)'][row + s_items.index("Solid Propellant")] = 0
                row += len(s_items)
        # multiply thicknesses by 0.001 m/mm after initializing .
        df['Thickness (m)'] = df['Thickness (m)']*0.001
    
    def initMasses(self, df, pl_items, l_items, s_items, load_type): # Divided additional engine mass of 59 kg by 3, wiring and avionics by 10

        # Load in max q data from matlab-generated .csv
        if load_type == 'Max-Q Condition':
                df_maxq = pd.read_csv('LVMasses\\Max Q Conditions_' + self.name + '.csv')
            

            #print(df_maxq)
            #print(df_maxq.iloc[0,3])
        # mass burned is on row 1, column 4 of df
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            if step.propulsion == 'Liquid':
                if i < len(self.listOfSteps) - 1: # if not last step
                    df['Mass (kg)'][row + l_items.index("Forward Skirt")] = step.interstage[1] * self.rho_body * df['Thickness (m)'][row + l_items.index("Forward Skirt")]
                
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['Mass (kg)'][row + l_items.index("Forward Skirt")] = step.fwd_skirt[1] * self.rho_body * df['Thickness (m)'][row + l_items.index("Forward Skirt")]

                df['Mass (kg)'][row + l_items.index("Avionics")] =  0.84 # mass of new avionics technology taken from NASA (https://technology.nasa.gov/patent/TOP2-274)
                df['Mass (kg)'][row + l_items.index("Wiring")] = 1.058*pow(self.m_0[i], 1/2) * pow(step.total_length, 1/4)/10
                # Skin Mass Calculation
                df['Mass (kg)'][row + l_items.index("Fuel Dome Top")] = step.dome_f[1] * step.rho_f_tank * df['Thickness (m)'][row + l_items.index("Fuel Dome Top")]
                df['Mass (kg)'][row + l_items.index("Fuel Cylinder")] = step.cyl_f[1] * step.rho_f_tank * df['Thickness (m)'][row + l_items.index("Fuel Cylinder")]
                df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")] = step.dome_f[1] * step.rho_f_tank * df['Thickness (m)'][row + l_items.index("Fuel Dome Top")]

                # MER Mass Calculation
                #tank_vol = 2 * step.dome_f[2] + step.cyl_f[2]
                #if step.propellants == 'Kerolox': 
                #    # relation for tank weight by mass (Table 8.6) for RP-1
                #    df['Mass (kg)'][row + l_items.index("Fuel Dome Top")] = step.dome_f[2] / tank_vol * 0.0148 * step.m_f_ideal 
                #    df['Mass (kg)'][row + l_items.index("Fuel Cylinder")] = step.cyl_f[2] / tank_vol * 0.0148 * step.m_f_ideal 
                #    df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")] = step.dome_f[2] / tank_vol * 0.0148 * step.m_f_ideal 
                #else:
                #    # relation for tank weight by volume (Table 8.7) for anything but LH2
                #    df['Mass (kg)'][row + l_items.index("Fuel Dome Top")] = step.dome_f[2] * 12.16
                #    df['Mass (kg)'][row + l_items.index("Fuel Cylinder")] = step.cyl_f[2] * 12.16
                #    df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")] = step.dome_f[2] * 12.16
                df['Mass (kg)'][row + l_items.index("Ox Dome Top")] = step.dome_ox[1] * step.rho_ox_tank * df['Thickness (m)'][row + l_items.index("Ox Dome Top")]
                df['Mass (kg)'][row + l_items.index("Ox Cylinder")] = step.cyl_ox[1] * step.rho_ox_tank * df['Thickness (m)'][row + l_items.index("Ox Cylinder")]
                df['Mass (kg)'][row + l_items.index("Ox Dome Bottom")] = step.dome_ox[1] * step.rho_ox_tank * df['Thickness (m)'][row + l_items.index("Ox Dome Bottom")]

                df['Mass (kg)'][row + l_items.index("Fuel Insulation")] = (step.cyl_f[1] + 2 * step.dome_f[1]) * step.SA_rho_insulation['Fuel']
                df['Mass (kg)'][row + l_items.index("Fuel Residual")] = step.residual_prop_perc * step.m_f_ideal
                df['Mass (kg)'][row + l_items.index("Intertank")] = step.intertank[1] * self.rho_body * df['Thickness (m)'][row + l_items.index("Intertank")]

                #tank_vol = 2 * step.dome_ox[2] + step.cyl_ox[2]
                # relation for tank weight by mass (Table 8.6) for LOX
                #df['Mass (kg)'][row + l_items.index("Ox Dome Top")] = step.dome_ox[2] / tank_vol * 0.0107 * step.m_ox_ideal
                #df['Mass (kg)'][row + l_items.index("Ox Cylinder")] = step.cyl_ox[2] / tank_vol * 0.0107 * step.m_ox_ideal
                #df['Mass (kg)'][row + l_items.index("Ox Dome Bottom")] = step.dome_ox[2] / tank_vol * 0.0107 * step.m_ox_ideal

                df['Mass (kg)'][row + l_items.index("Ox Insulation")] = (step.cyl_ox[1] + 2 * step.dome_ox[1]) * step.SA_rho_insulation['Oxidizer']  
                df['Mass (kg)'][row + l_items.index("Ox Residual")] = step.residual_prop_perc * step.m_ox_ideal
                df['Mass (kg)'][row + l_items.index("Pressurant Tank")] = step.press_tank[1] * self.rho_body * df['Thickness (m)'][row + l_items.index("Pressurant Tank")]
                #df['Mass (kg)'][row + l_items.index("Pressurant Tank")] = step.m_press * 2 # relation for tank weight by mass (Table 8.7)
                df['Mass (kg)'][row + l_items.index("Aft Skirt")] = step.aft_skirt[1] * self.rho_body *  df['Thickness (m)'][row + l_items.index("Aft Skirt")]  
                df['Mass (kg)'][row + l_items.index("Thrust Structure")] = 2.55*pow(10, -4)*step.T_SL
                df['Mass (kg)'][row + l_items.index("Gimballs")] = step.num_gimballed_engines * 237.8*pow(step.T_SL_engine / step.p_c, 0.9375)
                df['Mass (kg)'][row + l_items.index("Engines")] = step.num_engines * (step.T_SL_engine * (7.81 * pow(10, -4) + 3.37 * pow(10, -5) * sqrt(step.epsilon)) + self.m_engine_0)
                

                df['Mass (kg)'][row + l_items.index("Fuel")] = step.m_f_ideal + step.fuel_frac * step.startup_prop
                df['Mass (kg)'][row + l_items.index("Oxidizer")] = step.m_ox_ideal + step.ox_frac * step.startup_prop
                self.mp_actual.append(df['Mass (kg)'][row + l_items.index("Fuel")] + df['Mass (kg)'][row + l_items.index("Oxidizer")])
                if (load_type == "Max-Q Condition") & (step.step_num == 1):
                    df['Mass (kg)'][row + l_items.index("Fuel")] = df['Mass (kg)'][row + l_items.index("Fuel")] - df_maxq.iloc[0,4]*step.fuel_frac
                    df['Mass (kg)'][row + l_items.index("Oxidizer")] = df['Mass (kg)'][row + l_items.index("Oxidizer")] - df_maxq.iloc[0,4]*(1-step.fuel_frac)
                
                row += len(l_items)

            elif step.propulsion == 'Solid':
                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['Mass (kg)'][row + s_items.index("Nose Cone")] = (step.fairing[1] * self.rho_body * df['Thickness (m)'][row + s_items.index("Nose Cone")]) 
                        df['Mass (kg)'][row + s_items.index("Gimballs")] = 0 
                    else:
                        df['Mass (kg)'][row + s_items.index("Gimballs")] = step.num_gimballed_engines * 237.8*pow(step.T_SL_engine / step.p_c, 0.9375)
                    df['Mass (kg)'][row + s_items.index("Forward Skirt")] = step.interstage[1] * self.rho_body * df['Thickness (m)'][row + s_items.index("Forward Skirt")] 
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['Mass (kg)'][row + s_items.index("Forward Skirt")] = step.fwd_skirt[1] * self.rho_body * df['Thickness (m)'][row + s_items.index("Forward Skirt")] 

                df['Mass (kg)'][row + s_items.index("Avionics")] =  0.84 # mass of new avionics technology taken from NASA (https://technology.nasa.gov/patent/TOP2-274)
                df['Mass (kg)'][row + s_items.index("Wiring")] = 1.058*pow(self.m_0[i], 1/2) * pow(step.total_length, 1/4)/10 
                df['Mass (kg)'][row + s_items.index("Pressurant Tank")] = step.press_tank[1] * self.rho_body * df['Thickness (m)'][row + s_items.index("Pressurant Tank")]
                #df['Mass (kg)'][row + s_items.index("Pressurant Tank")] = step.m_press * 2 # relation for tank weight (Table 8.7)
                m_srm = step.k_SRM * step.m_prop_tot
                SA_srm = 2*step.dome_f[1]+step.srm_casing[1]
                #m_srm = SA_srm*self.rho_body*0.001
                df['Mass (kg)'][row + s_items.index("SRM Dome Top")] = m_srm * step.dome_f[1] / SA_srm # mass by percent SA of SRM
                df['Mass (kg)'][row + s_items.index("Solid Propellant Casing")] = m_srm * step.srm_casing[1] / SA_srm # mass by percent SA of SRM
                df['Mass (kg)'][row + s_items.index("SRM Dome Bottom")] = m_srm * step.dome_f[1] / SA_srm # mass by percent SA of SRM
                df['Mass (kg)'][row + s_items.index("Solid Propellant Residual")] = step.residual_prop
                df['Mass (kg)'][row + s_items.index("Aft Skirt")] = step.aft_skirt[1] * self.rho_body *  df['Thickness (m)'][row + s_items.index("Aft Skirt")]
                df['Mass (kg)'][row + s_items.index("Nozzle")] = 0
                df['Mass (kg)'][row + s_items.index("Solid Propellant")] = step.m_prop_tot
                self.mp_actual.append(step.m_prop_tot*step.multiplier)
                if (load_type == "Max-Q Condition"):
                    if (self.name == 'Latona-1'):
                        if step.step_num == 1:
                            df['Mass (kg)'][row + s_items.index("Solid Propellant")] = df['Mass (kg)'][row + s_items.index("Solid Propellant")] - df_maxq.iloc[0][4] # Latona 1 case
                    elif self.name == 'Latona-2':
                        if (step.step_num == 1) & step.parallel:
                            print(df_maxq.iloc[0,5])
                            df['Mass (kg)'][row + s_items.index("Solid Propellant")] = df_maxq.iloc[0,5] # propellant mass OF EACH booster
                        elif step.step_num == 2:
                            print(df_maxq.iloc[0,6])
                            df['Mass (kg)'][row + s_items.index("Solid Propellant")] = df_maxq.iloc[0,6] # propellant mass of main stage

                row += len(s_items)

            df['Mass (kg)'] = df['Mass (kg)']*step.multiplier # be careful here multiplying the mass column for each step iteration. This will result in 
            # multiplying all the previously instantiated step's masses in the dataframe by the current's step's multiplier. 
            # It's okay here since only the 1st step's multiplier != 1, but be aware.
            # could fix by entering inside loop and doing if step.parallel -> if step.propulsion == liq or sol -> 
            # df['Mass (kg)'][index_start:index_end]=df['Mass (kg)'][index_start:index_end]*step.multiplier, where index_end is the current row + s/l_items.index and 
            # index_start = index_end - (num_l/s items)
            df['Mass (kg)'][pl_items.index("PLF")] = sum(self.listOfSteps[len(self.listOfSteps)-1].fairing[1]) * self.listOfSteps[len(self.listOfSteps)-1].rho_fairing * df['Thickness (m)'][pl_items.index("PLF")]
            #print("Payload Fairing Cylinder SA is " + str(self.listOfSteps[len(self.listOfSteps)-1].fairing[1][0]))
            #print("Payload Fairing Cone SA is " + str(self.listOfSteps[len(self.listOfSteps)-1].fairing[1][1]))
            #print("Total Payload Fairing SA is " + str(sum(self.listOfSteps[len(self.listOfSteps)-1].fairing[1])))
            df['Mass (kg)'][pl_items.index("Payload")] = self.PL                                  
            df['Mass (kg)'][pl_items.index("PAF")] = self.m_PAF
    
    def initDistances(self, df, pl_items, l_items, s_items, load_type): # initiate distances from bottom of aft skirt to CM of component
        # Note the distances are initialized in reverse order, but maintain their order in the dataframe df
        pi = 3.1415926535897932
        row = len(pl_items)
        CG = 0
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            if step.propulsion == "Liquid":
                CG += step.r/4
                df['Distance (m)'][row + l_items.index("Engines")] = CG
                CG += step.L_n - step.r/4
                df['Distance (m)'][row + l_items.index("Gimballs")] = CG     
                CG += step.T_struct/2 - step.L_n
                df['Distance (m)'][row + l_items.index("Thrust Structure")] = CG           
                CG += step.aft_skirt[0]/2 - step.T_struct/2
                df['Distance (m)'][row + l_items.index("Aft Skirt")] = CG
                CG += step.aft_skirt[0]/2 - 4 * step.dome_ox[0]/(3 * pi) #+ step.dome_ox[0] # add to CG of ox dome bottom (gap of step.r/2)
                print("CG of dome from bottom of dome = " + str(4 * step.dome_ox[0]/(3 * pi)))
                df['Distance (m)'][row + l_items.index("Ox Dome Bottom")] = CG   
                CG += 4 * step.dome_ox[0]/(3 * pi) + step.cyl_ox[0]/2
                df['Distance (m)'][row + l_items.index("Ox Cylinder")] = CG 
                df['Distance (m)'][row + l_items.index("Ox Insulation")] = CG  
                df['Distance (m)'][row + l_items.index("Ox Residual")] = CG - step.cyl_ox[0]/2 - step.dome_ox[0]/2 # residual CG sits in the middle of the bottom ox dome tank
                df['Distance (m)'][row + l_items.index("Pressurant Tank")] = CG
                # if load_type == "Ground Wind-Loads Condition":
                #     df['Distance (m)'][row + l_items.index("Oxidizer")] = CG
                # elif load_type == "Max-Q Condition":
                #    # max q condition. delete for ground load condition
                vol_ox = df['Mass (kg)'][row + l_items.index("Oxidizer")]/step.rho_ox
                vol_ox_cyl_full = vol_ox - step.dome_ox[2]
                m_dome_ox = step.dome_ox[2] * step.rho_ox
                m_cyl_ox = vol_ox_cyl_full * step.rho_ox
                cg_dome = step.dome_ox[0]- 4 * step.dome_ox[0]/(3 * pi) # from bottom of dome
                cg_cyl = step.dome_ox[0] + (vol_ox_cyl_full/(pi * pow(step.r,2)))/2 # from bottom of dome
                cg_ox = (m_dome_ox * cg_dome + m_cyl_ox * cg_cyl)/(m_dome_ox + m_cyl_ox) # resultant cg of oxidizer from bottom of dome
                df['Distance (m)'][row + l_items.index("Oxidizer")] = CG - step.cyl_ox[0]/2 - step.dome_ox[0] + cg_ox # subtract cg back down to bottom of ox dome and add cg_ox from there
                
                #df['Distance (m)'][row + l_items.index("Oxidizer")] = CG
                CG += step.cyl_ox[0]/2 + 4 * step.dome_ox[0]/(3 * pi)
                df['Distance (m)'][row + l_items.index("Ox Dome Top")] = CG
                CG += - 4 * step.dome_ox[0]/(3 * pi) + step.intertank[0]/2
                df['Distance (m)'][row + l_items.index("Intertank")] = CG 
                CG += step.r/4 + step.dome_f[0] - 4 * step.dome_f[0]/(3 * pi)
                df['Distance (m)'][row + l_items.index("Fuel Dome Bottom")] = CG
                CG += 4 * step.dome_f[0]/(3 * pi) + step.cyl_f[0]/2
                df['Distance (m)'][row + l_items.index("Fuel Cylinder")] = CG
                df['Distance (m)'][row + l_items.index("Fuel Insulation")] = CG
                df['Distance (m)'][row + l_items.index("Fuel Residual")] = CG - step.cyl_f[0]/2 - step.dome_f[0]/2 # residual CG sits in the middle of the bottom fuel dome tank
                # if load_type == "Ground Wind-Loads Condition":
                #     df['Distance (m)'][row + l_items.index("Fuel")] = CG
                # elif load_type == "Max-Q Condition":
                vol_f = df['Mass (kg)'][row + l_items.index("Fuel")]/step.rho_f
                vol_f_cyl_full = vol_f - step.dome_f[2]
                m_dome_f = step.dome_f[2] * step.rho_f
                m_cyl_f = vol_f_cyl_full * step.rho_f
                CG_dome = step.dome_f[0]- 4 * step.dome_f[0]/(3 * pi) # from bottom of dome
                CG_cyl = step.dome_f[0] + (vol_f_cyl_full/(pi * pow(step.r,2)))/2 # from bottom of dome
                CG_f = (m_dome_f * CG_dome + m_cyl_f * CG_cyl)/(m_dome_f + m_cyl_f) # Resultant CG of fuel from bottom of dome
                df['Distance (m)'][row + l_items.index("Fuel")] = CG - step.cyl_f[0]/2 - step.dome_f[0] + CG_f # subtract CG back down to bottom of fuel dome and add CG_f from there

                CG += step.cyl_f[0]/2 + 4 * step.dome_f[0]/(3 * pi)
                df['Distance (m)'][row + l_items.index("Fuel Dome Top")] = CG
                df['Distance (m)'][row + l_items.index("Wiring")] =  0

                if i < len(self.listOfSteps) - 1: # if not last step NOTE: invert step.interstage CG calc by getting rid of step.interstage[0] and changing the '-' (CG calc) to a '+'                    CG += - 4 * step.dome_f[0]/(3 * pi) + step.interstage[0] - step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r)  #formual for centroid of trapezoid)
                    CG += -4 * step.dome_f[0]/(3 * pi) + step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r)
                    df['Distance (m)'][row + l_items.index("Forward Skirt")] = CG
                    df['Distance (m)'][row + l_items.index("Avionics")] = CG
                    CG += step.interstage[0] - step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r)
                elif i == len(self.listOfSteps) - 1: # if last step
                    CG +=  step.fwd_skirt[0]/2 - 4 * step.dome_f[0]/(3 * pi) 
                    df['Distance (m)'][row + l_items.index("Forward Skirt")] = CG
                    df['Distance (m)'][row + l_items.index("Avionics")] = CG
                    CG += step.fwd_skirt[0]/2
                    if self.Mission.input[0] == 'One':
                        df['Distance (m)'][pl_items.index("PAF")] = CG # Mission 1: place distance of PAF at bottom of PLF cylinder
                    elif self.Mission.input[0] == 'Two':
                        df['Distance (m)'][pl_items.index("PAF")] = CG + step.fairing[0][0]/2 # Mission 2: place distance of PAF in the middle of PLF cylinder 
                    CG += step.fairing[0][0]/2 
                    df['Distance (m)'][pl_items.index("Payload")] = CG # middle of fairing_cylinder
                    CG -= step.fairing[0][0]/2 # Subtract CG to get back down to PAF
                    # CG of payload fairing = (m_cyl*CG_cyl + m_cone*CG_cone)/m_total  
                    CG += (step.fairing[1][0]*step.rho_fairing*df['Thickness (m)'][pl_items.index("PLF")]*step.fairing[0][0]/2 + step.fairing[1][1]*step.rho_fairing*df['Thickness (m)'][pl_items.index("PLF")]*(step.fairing[0][1]/3 + step.fairing[0][0])) / (df['Mass (kg)'][pl_items.index("PLF")]) # CG of payload fairing = (m_cyl*CG_cyl + m_cone*CG_cone)/m_total        (CG from PAF)  
                    df['Distance (m)'][pl_items.index("PLF")] = CG
                    # Subtract CG to get back down to PAF
                    CG -= (step.fairing[1][0]*step.rho_fairing*df['Thickness (m)'][pl_items.index("PLF")]*step.fairing[0][0]/2 + step.fairing[1][1]*step.rho_fairing*df['Thickness (m)'][pl_items.index("PLF")]*(step.fairing[0][1]/3 + step.fairing[0][0])) / (df['Mass (kg)'][pl_items.index("PLF")]) # CG of payload fairing = (m_cyl*CG_cyl + m_cone*CG_cone)/m_total        (CG from PAF)
                    CG += sum(step.fairing[0]) # add total fairing height on top of PAF to get top of PLF
                    df['Distance (m)'][0] = CG  # Top of PLF
                
                row += len(l_items)

            elif step.propulsion == 'Solid':
                CG += step.r/4
                df['Distance (m)'][row + s_items.index("Nozzle")] = CG 
                CG += step.L_n - step.r/4
                df['Distance (m)'][row + s_items.index("Gimballs")] = CG 
                CG += step.aft_skirt[0]/2 - step.L_n
                df['Distance (m)'][row + s_items.index("Aft Skirt")] = CG 
                CG += step.aft_skirt[0]/2 + 4 * step.dome_f[0]/(3 * pi)
                df['Distance (m)'][row + s_items.index("SRM Dome Bottom")] = CG 
                CG += -4 * step.dome_f[0]/(3 * pi) + step.srm_casing[0]/2
                df['Distance (m)'][row + s_items.index("Solid Propellant Casing")] = CG 
                df['Distance (m)'][row + s_items.index("Solid Propellant Residual")] = CG
                if load_type == "Ground Wind-Loads Condition":
                    df['Distance (m)'][row + s_items.index("Solid Propellant")] = CG 
                elif load_type == "Max-Q Condition":
                    df['Distance (m)'][row + s_items.index("Solid Propellant")] = CG  # Same CG since the solid propellant burns out radially
                
                CG += step.srm_casing[0]/2 + 4 * step.dome_f[0]/(3 * pi)
                df['Distance (m)'][row + s_items.index("SRM Dome Top")] = CG 
                df['Distance (m)'][row + s_items.index("Wiring")] = 0
                
                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel: # interstage is cylindrical;
                        #start from bottom of SRM Dome Top
                        CG += -4 * step.dome_f[0]/(3 * pi) + step.interstage[0]/2 
                        df['Distance (m)'][row + s_items.index("Forward Skirt")] = CG
                        df['Distance (m)'][row + s_items.index("Avionics")] = CG 
                        CG += -step.interstage[0]/2 + step.r/2 + step.press_tank[0] # add gap of r/2 between solid propellant casing and pressurant tank
                        df['Distance (m)'][row + s_items.index("Pressurant Tank")] = CG
                        CG += -step.r/2 - step.press_tank[0] + step.interstage[0] + step.fairing[0]/3
                        df['Distance (m)'][row + s_items.index("Nose Cone")] = CG
                    else: # interstage is a cut-off-cone
                        #start from SRM Dome Top CG
                        CG += -4 * step.dome_f[0]/(3 * pi) + step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r) # ( using formula for centroid of trapezoid)
                        df['Distance (m)'][row + s_items.index("Forward Skirt")] = CG
                        fwd_skirt_CG = CG
                        df['Distance (m)'][row + s_items.index("Avionics")] = CG 
                        CG += -step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r) + step.r/2 + step.press_tank[0] # add gap of r/2 between solid propellant casing and pressurant tank
                        df['Distance (m)'][row + s_items.index("Pressurant Tank")] = CG
                        CG = fwd_skirt_CG + step.interstage[0] - step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r)
       
                elif i == len(self.listOfSteps) - 1: # if last step
                    CG += -4 * step.dome_f[0]/(3 * pi)  # set CG to bottom of SRM Dome Top
                    df['Distance (m)'][row + s_items.index("Pressurant Tank")] = CG + 4 * step.dome_f[0]/(3 * pi) + step.press_tank[0] + step.r/2 # gap of step.r/2 from top of SRM Dome Top to bottom of pressure tank
                    CG += step.fwd_skirt[0]/2 # add to CG of fwd skirt starting from bottom of SRM Dome Top
                    df['Distance (m)'][row + s_items.index("Forward Skirt")] = CG
                    df['Distance (m)'][row + s_items.index("Avionics")] = CG 
                    CG += step.fwd_skirt[0]/2
                    df['Distance (m)'][pl_items.index("PAF")] = CG
                    CG += step.fairing[0][0]/2
                    df['Distance (m)'][pl_items.index("Payload")] = CG
                    CG -= step.fairing[0][0]/2
                    #step.fairing[1][0]*self.rho_fairing
                    df['Distance (m)'][pl_items.index("PLF")] = CG 
                row += len(s_items)
            if step.parallel: # be careful here, only properly sets CG if parallel is the 1st stage, not parallel subsequent stages (although it would never occur otherwise, be aware)
                CG = 0

    def initMoments(self, df, pl_items, l_items, s_items): # initiate distances from bottom of aft skirt to CM of component
        # Note the distances are initialized in reverse order, but maintain their order in the dataframe df
        #pi = 3.1415926535897932
        row = len(pl_items)
        CG = 0
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            
            if step.propulsion == "Liquid":
                df['Moment (kg*m)'][row + l_items.index("Avionics")] = df['Distance (m)'][row + l_items.index("Avionics")] * df['Mass (kg)'][row + l_items.index("Avionics")]    
                df['Moment (kg*m)'][row + l_items.index("Wiring")] =  0          
                df['Moment (kg*m)'][row + l_items.index("Fuel Dome Top")] = df['Distance (m)'][row + l_items.index("Fuel Dome Top")] * df['Mass (kg)'][row + l_items.index("Fuel Dome Top")]  
                df['Moment (kg*m)'][row + l_items.index("Fuel Cylinder")] = df['Distance (m)'][row + l_items.index("Fuel Cylinder")] * df['Mass (kg)'][row + l_items.index("Fuel Cylinder")]       
                df['Moment (kg*m)'][row + l_items.index("Fuel Dome Bottom")] = df['Distance (m)'][row + l_items.index("Fuel Dome Bottom")] * df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")]  
                df['Moment (kg*m)'][row + l_items.index("Fuel Insulation")] = df['Distance (m)'][row + l_items.index("Fuel Insulation")] * df['Mass (kg)'][row + l_items.index("Fuel Insulation")]  
                df['Moment (kg*m)'][row + l_items.index("Fuel Residual")] = df['Distance (m)'][row + l_items.index("Fuel Residual")] * df['Mass (kg)'][row + l_items.index("Fuel Residual")] 
                df['Moment (kg*m)'][row + l_items.index("Intertank")] = df['Distance (m)'][row + l_items.index("Intertank")] * df['Mass (kg)'][row + l_items.index("Intertank")]   
                df['Moment (kg*m)'][row + l_items.index("Ox Dome Top")] = df['Distance (m)'][row + l_items.index("Ox Dome Top")] * df['Mass (kg)'][row + l_items.index("Ox Dome Top")]      
                df['Moment (kg*m)'][row + l_items.index("Ox Cylinder")] = df['Distance (m)'][row + l_items.index("Ox Cylinder")] * df['Mass (kg)'][row + l_items.index("Ox Cylinder")]    
                df['Moment (kg*m)'][row + l_items.index("Ox Dome Bottom")] = df['Distance (m)'][row + l_items.index("Ox Dome Bottom")] * df['Mass (kg)'][row + l_items.index("Ox Dome Bottom")]   
                df['Moment (kg*m)'][row + l_items.index("Ox Insulation")] = df['Distance (m)'][row + l_items.index("Ox Insulation")] * df['Mass (kg)'][row + l_items.index("Ox Insulation")]  
                df['Moment (kg*m)'][row + l_items.index("Ox Residual")] = df['Distance (m)'][row + l_items.index("Ox Residual")] * df['Mass (kg)'][row + l_items.index("Ox Residual")]  
                df['Moment (kg*m)'][row + l_items.index("Pressurant Tank")] = df['Distance (m)'][row + l_items.index("Pressurant Tank")] * df['Mass (kg)'][row + l_items.index("Pressurant Tank")]
                df['Moment (kg*m)'][row + l_items.index("Aft Skirt")] = df['Distance (m)'][row + l_items.index("Aft Skirt")] * df['Mass (kg)'][row + l_items.index("Aft Skirt")]   
                df['Moment (kg*m)'][row + l_items.index("Thrust Structure")] = df['Distance (m)'][row + l_items.index("Thrust Structure")] * df['Mass (kg)'][row + l_items.index("Thrust Structure")]  
                df['Moment (kg*m)'][row + l_items.index("Gimballs")] = df['Distance (m)'][row + l_items.index("Gimballs")] * df['Mass (kg)'][row + l_items.index("Gimballs")] 
                df['Moment (kg*m)'][row + l_items.index("Engines")] = df['Distance (m)'][row + l_items.index("Engines")] * df['Mass (kg)'][row + l_items.index("Engines")] 
                df['Moment (kg*m)'][row + l_items.index("Fuel")] = df['Distance (m)'][row + l_items.index("Fuel")] * df['Mass (kg)'][row + l_items.index("Fuel")]   
                df['Moment (kg*m)'][row + l_items.index("Oxidizer")] = df['Distance (m)'][row + l_items.index("Oxidizer")] * df['Mass (kg)'][row + l_items.index("Oxidizer")]  

                if i < len(self.listOfSteps) - 1: # if not last step
                    CG += step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r) - 4 * step.dome_f[0]/(3 * pi) #( using formual for centroid of trapezoid)
                    df['Moment (kg*m)'][row + l_items.index("Forward Skirt")] = df['Distance (m)'][row + l_items.index("Forward Skirt")] * df['Mass (kg)'][row + l_items.index("Forward Skirt")]
                    CG += step.interstage[0] - step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r)
                elif i == len(self.listOfSteps) - 1: # if last step
                    CG +=  step.fwd_skirt[0]/2 - 4 * step.dome_f[0]/(3 * pi) # fwd skirt
                    df['Moment (kg*m)'][row + l_items.index("Forward Skirt")] = df['Distance (m)'][row + l_items.index("Forward Skirt")] * df['Mass (kg)'][row + l_items.index("Forward Skirt")]                # Interstage/Forward Skirt
                    df['Moment (kg*m)'][pl_items.index("PAF")] = df['Distance (m)'][pl_items.index("PAF")] * df['Mass (kg)'][pl_items.index("PAF")]                  
                    df['Moment (kg*m)'][pl_items.index("Payload")] = df['Distance (m)'][pl_items.index("Payload")] * df['Mass (kg)'][pl_items.index("Payload")]                  
                    df['Moment (kg*m)'][pl_items.index("PLF")] = df['Distance (m)'][pl_items.index("PLF")] * df['Mass (kg)'][pl_items.index("PLF")]                 
                
                row += len(l_items)

            elif step.propulsion == 'Solid':

                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['Moment (kg*m)'][row + s_items.index("Nose Cone")] = df['Distance (m)'][row + s_items.index("Nose Cone")] * df['Mass (kg)'][row + s_items.index("Nose Cone")] 
                    df['Moment (kg*m)'][row + s_items.index("Forward Skirt")] = df['Distance (m)'][row + s_items.index("Forward Skirt")] * df['Mass (kg)'][row + s_items.index("Forward Skirt")]
                elif i == len(self.listOfSteps) - 1: # if last step     
                    df['Moment (kg*m)'][row + s_items.index("Forward Skirt")] = df['Distance (m)'][row + s_items.index("Forward Skirt")] * df['Mass (kg)'][row + s_items.index("Forward Skirt")]
                    df['Moment (kg*m)'][pl_items.index("PAF")] = df['Distance (m)'][pl_items.index("PAF")] * df['Mass (kg)'][pl_items.index("PAF")]                  
                    df['Moment (kg*m)'][pl_items.index("Payload")] = df['Distance (m)'][pl_items.index("Payload")] * df['Mass (kg)'][pl_items.index("Payload")]     
                    df['Moment (kg*m)'][pl_items.index("PLF")] = df['Distance (m)'][pl_items.index("PLF")] * df['Mass (kg)'][pl_items.index("PLF")]     

                df['Moment (kg*m)'][row + s_items.index("Avionics")] = df['Distance (m)'][row + s_items.index("Avionics")] * df['Mass (kg)'][row + s_items.index("Avionics")]  
                df['Moment (kg*m)'][row + s_items.index("Wiring")] = 0   
                df['Moment (kg*m)'][row + s_items.index("Pressurant Tank")] = df['Distance (m)'][row + s_items.index("Pressurant Tank")] * df['Mass (kg)'][row + s_items.index("Pressurant Tank")] 
                df['Moment (kg*m)'][row + s_items.index("SRM Dome Top")] = df['Distance (m)'][row + s_items.index("SRM Dome Top")] * df['Mass (kg)'][row + s_items.index("SRM Dome Top")]   
                df['Moment (kg*m)'][row + s_items.index("Solid Propellant Casing")] = df['Distance (m)'][row + s_items.index("Solid Propellant Casing")] * df['Mass (kg)'][row + s_items.index("Solid Propellant Casing")]   
                df['Moment (kg*m)'][row + s_items.index("SRM Dome Bottom")] = df['Distance (m)'][row + s_items.index("SRM Dome Bottom")] * df['Mass (kg)'][row + s_items.index("SRM Dome Bottom")]   
                df['Moment (kg*m)'][row + s_items.index("Solid Propellant Residual")] = df['Distance (m)'][row + s_items.index("Solid Propellant Residual")] * df['Mass (kg)'][row + s_items.index("Solid Propellant Residual")]               
                df['Moment (kg*m)'][row + s_items.index("Aft Skirt")] = df['Distance (m)'][row + s_items.index("Aft Skirt")] * df['Mass (kg)'][row + s_items.index("Aft Skirt")]  
                df['Moment (kg*m)'][row + s_items.index("Gimballs")] = df['Distance (m)'][row + s_items.index("Gimballs")] * df['Mass (kg)'][row + s_items.index("Gimballs")]
                df['Moment (kg*m)'][row + s_items.index("Nozzle")] = df['Distance (m)'][row + s_items.index("Nozzle")] * df['Mass (kg)'][row + s_items.index("Nozzle")]  
                df['Moment (kg*m)'][row + s_items.index("Solid Propellant")] = df['Distance (m)'][row + s_items.index("Solid Propellant")] * df['Mass (kg)'][row + s_items.index("Solid Propellant")] 
                row += len(s_items)

    def initDistFromCM(self, df, pl_items, l_items, s_items): # initiate distances from bottom of aft skirt to CM of component
        # Note the distances are initialized in reverse order, but maintain their order in the dataframe df
        #pi = 3.1415926535897932
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            
            if step.propulsion == "Liquid":
                df['Distance from CM (m)'][row + l_items.index("Avionics")] = df['Distance (m)'][row + l_items.index("Avionics")] - self.CM_full 
                df['Distance from CM (m)'][row + l_items.index("Wiring")] =  0              
                df['Distance from CM (m)'][row + l_items.index("Fuel Dome Top")] = df['Distance (m)'][row + l_items.index("Fuel Dome Top")] - self.CM_full      
                df['Distance from CM (m)'][row + l_items.index("Fuel Cylinder")] = df['Distance (m)'][row + l_items.index("Fuel Cylinder")] - self.CM_full         
                df['Distance from CM (m)'][row + l_items.index("Fuel Dome Bottom")] = df['Distance (m)'][row + l_items.index("Fuel Dome Bottom")] - self.CM_full    
                df['Distance from CM (m)'][row + l_items.index("Fuel Insulation")] = df['Distance (m)'][row + l_items.index("Fuel Insulation")] - self.CM_full     
                df['Distance from CM (m)'][row + l_items.index("Fuel Residual")] = df['Distance (m)'][row + l_items.index("Fuel Residual")] - self.CM_full                 
                df['Distance from CM (m)'][row + l_items.index("Intertank")] = df['Distance (m)'][row + l_items.index("Intertank")] - self.CM_full         
                df['Distance from CM (m)'][row + l_items.index("Ox Dome Top")] = df['Distance (m)'][row + l_items.index("Ox Dome Top")] - self.CM_full    
                df['Distance from CM (m)'][row + l_items.index("Ox Cylinder")] = df['Distance (m)'][row + l_items.index("Ox Cylinder")] - self.CM_full       
                df['Distance from CM (m)'][row + l_items.index("Ox Dome Bottom")] = df['Distance (m)'][row + l_items.index("Ox Dome Bottom")] - self.CM_full   
                df['Distance from CM (m)'][row + l_items.index("Ox Insulation")] = df['Distance (m)'][row + l_items.index("Ox Insulation")] - self.CM_full   
                df['Distance from CM (m)'][row + l_items.index("Ox Residual")] = df['Distance (m)'][row + l_items.index("Ox Residual")] - self.CM_full   
                df['Distance from CM (m)'][row + l_items.index("Pressurant Tank")] = df['Distance (m)'][row + l_items.index("Pressurant Tank")] - self.CM_full
                df['Distance from CM (m)'][row + l_items.index("Aft Skirt")] = df['Distance (m)'][row + l_items.index("Aft Skirt")] - self.CM_full  
                df['Distance from CM (m)'][row + l_items.index("Thrust Structure")] = df['Distance (m)'][row + l_items.index("Thrust Structure")] - self.CM_full  
                df['Distance from CM (m)'][row + l_items.index("Gimballs")] = df['Distance (m)'][row + l_items.index("Gimballs")] - self.CM_full  
                df['Distance from CM (m)'][row + l_items.index("Engines")] = df['Distance (m)'][row + l_items.index("Engines")] - self.CM_full  
                df['Distance from CM (m)'][row + l_items.index("Fuel")] = df['Distance (m)'][row + l_items.index("Fuel")] - self.CM_full  
                df['Distance from CM (m)'][row + l_items.index("Oxidizer")] = df['Distance (m)'][row + l_items.index("Oxidizer")] - self.CM_full  

                if i < len(self.listOfSteps) - 1: # if not last step
                    # Interstage ( using formual for centroid of trapezoid)
                    df['Distance from CM (m)'][row + l_items.index("Forward Skirt")] = df['Distance (m)'][row + l_items.index("Forward Skirt")] - self.CM_full
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['Distance from CM (m)'][row + l_items.index("Forward Skirt")] = df['Distance (m)'][row + l_items.index("Forward Skirt")] - self.CM_full             
                    df['Distance from CM (m)'][pl_items.index("PAF")] = df['Distance (m)'][pl_items.index("PAF")] - self.CM_full                  
                    df['Distance from CM (m)'][pl_items.index("Payload")] = df['Distance (m)'][pl_items.index("Payload")] - self.CM_full                  
                    df['Distance from CM (m)'][pl_items.index("PLF")] = df['Distance (m)'][pl_items.index("PLF")] - self.CM_full                 
                
                row += len(l_items)

            elif step.propulsion == 'Solid':

                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['Distance from CM (m)'][row + s_items.index("Nose Cone")] = df['Distance (m)'][row + s_items.index("Nose Cone")] - self.CM_full 
                    df['Distance from CM (m)'][row + s_items.index("Forward Skirt")] = df['Distance (m)'][row + s_items.index("Forward Skirt")] - self.CM_full
                elif i == len(self.listOfSteps) - 1: # if last step          
                    df['Distance from CM (m)'][row + s_items.index("Forward Skirt")] = df['Distance (m)'][row + s_items.index("Forward Skirt")] - self.CM_full
                    df['Distance from CM (m)'][pl_items.index("PAF")] = df['Distance (m)'][pl_items.index("PAF")] - self.CM_full                 
                    df['Distance from CM (m)'][pl_items.index("Payload")] = df['Distance (m)'][pl_items.index("Payload")] - self.CM_full                 
                    df['Distance from CM (m)'][pl_items.index("PLF")] = df['Distance (m)'][pl_items.index("PLF")] - self.CM_full         

                df['Distance from CM (m)'][row + s_items.index("Avionics")] = df['Distance (m)'][row + s_items.index("Avionics")] - self.CM_full   
                df['Distance from CM (m)'][row + s_items.index("Wiring")] = 0  
                df['Distance from CM (m)'][row + s_items.index("Pressurant Tank")] = df['Distance (m)'][row + s_items.index("Pressurant Tank")] - self.CM_full  
                df['Distance from CM (m)'][row + s_items.index("SRM Dome Top")] = df['Distance (m)'][row + s_items.index("SRM Dome Top")] - self.CM_full  
                df['Distance from CM (m)'][row + s_items.index("Solid Propellant Casing")] = df['Distance (m)'][row + s_items.index("Solid Propellant Casing")] - self.CM_full  
                df['Distance from CM (m)'][row + s_items.index("SRM Dome Bottom")] = df['Distance (m)'][row + s_items.index("SRM Dome Bottom")] - self.CM_full  
                df['Distance from CM (m)'][row + s_items.index("Solid Propellant Residual")] = 0 
                df['Distance from CM (m)'][row + s_items.index("Aft Skirt")] = df['Distance (m)'][row + s_items.index("Aft Skirt")] - self.CM_full
                df['Distance from CM (m)'][row + s_items.index("Gimballs")] = df['Distance (m)'][row + s_items.index("Gimballs")] - self.CM_full
                df['Distance from CM (m)'][row + s_items.index("Nozzle")] = df['Distance (m)'][row + s_items.index("Nozzle")] - self.CM_full
                df['Distance from CM (m)'][row + s_items.index("Solid Propellant")] = df['Distance (m)'][row + s_items.index("Solid Propellant")] - self.CM_full
                row += len(s_items)
    
    # init J0's - double check some values such as payload, fairing, propellant tanks etc
    def initJ0sOld(self, df, pl_items, l_items, s_items): 
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            
            if step.propulsion == "Liquid":
                df['J0 (kg m^2)'][row + l_items.index("Avionics")] = df['Mass (kg)'][row + l_items.index("Avionics")] * pow(step.r,2)/2     
                df['J0 (kg m^2)'][row + l_items.index("Wiring")] = df['Mass (kg)'][row + l_items.index("Wiring")] * pow(step.r,2)/2               
                df['J0 (kg m^2)'][row + l_items.index("Fuel Dome Top")] = df['Mass (kg)'][row + l_items.index("Fuel Dome Top")] * pow(step.r,2)     
                df['J0 (kg m^2)'][row + l_items.index("Fuel Cylinder")] = J0ThinCyl(df['Mass (kg)'][row + l_items.index("Fuel Cylinder")], step.r, step.cyl_f[0])
                df['J0 (kg m^2)'][row + l_items.index("Fuel Dome Bottom")] = df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")] * pow(step.r,2)     
                df['J0 (kg m^2)'][row + l_items.index("Fuel Insulation")] = df['Mass (kg)'][row + l_items.index("Fuel Insulation")] * (pow(step.r,2)/2 + pow(df['Height (m)'][row + l_items.index("Fuel Insulation")],2)/12)      
                df['J0 (kg m^2)'][row + l_items.index("Fuel Residual")] = 0   
                df['J0 (kg m^2)'][row + l_items.index("Intertank")] = J0ThinCyl(df['Mass (kg)'][row + l_items.index("Intertank")], step.r, step.intertank[0])
                df['J0 (kg m^2)'][row + l_items.index("Ox Dome Top")] = df['Mass (kg)'][row + l_items.index("Ox Dome Top")] * pow(step.r,2)     
                df['J0 (kg m^2)'][row + l_items.index("Ox Cylinder")] = J0ThinCyl(df['Mass (kg)'][row + l_items.index("Intertank")], step.r, step.intertank[0])
                df['J0 (kg m^2)'][row + l_items.index("Ox Dome Bottom")] = df['Mass (kg)'][row + l_items.index("Ox Dome Bottom")] * pow(step.r,2)
                df['J0 (kg m^2)'][row + l_items.index("Ox Insulation")] = df['Mass (kg)'][row + l_items.index("Ox Insulation")] * (pow(step.r,2)/2 + pow(df['Height (m)'][row + l_items.index("Ox Insulation")],2)/12)   
                df['J0 (kg m^2)'][row + l_items.index("Ox Residual")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Pressurant Tank")] = df['Mass (kg)'][row + l_items.index("Pressurant Tank")] * pow(step.r,2)
                df['J0 (kg m^2)'][row + l_items.index("Aft Skirt")] = df['Mass (kg)'][row + l_items.index("Aft Skirt")] * ( pow(step.r,2) + pow(df['Height (m)'][row + l_items.index("Aft Skirt")],2)/12) 
                df['J0 (kg m^2)'][row + l_items.index("Thrust Structure")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Gimballs")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Engines")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Fuel")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Oxidizer")] = 0

                if i < len(self.listOfSteps) - 1: # if not last step
                    # Interstage ( using formual for centroid of trapezoid)
                    df['J0 (kg m^2)'][row + l_items.index("Forward Skirt")] = df['Mass (kg)'][row + l_items.index("Forward Skirt")] * ( (pow(step.r,2) + pow(self.listOfSteps[i+1].r,2))/4 + pow(df['Height (m)'][row + l_items.index("Forward Skirt")],2)/18 * (1+(2*step.r*self.listOfSteps[i+1].r)/(pow(step.r + self.listOfSteps[i+1].r,2))))
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['J0 (kg m^2)'][row + l_items.index("Forward Skirt")] = df['Mass (kg)'][row + l_items.index("Forward Skirt")] * (pow(step.r,2)/2 + pow(df['Height (m)'][row + l_items.index("Forward Skirt")],2)/12)          
                    df['J0 (kg m^2)'][pl_items.index("PAF")] = df['Mass (kg)'][pl_items.index("PAF")] * pow(step.r,2)/2                  
                    df['J0 (kg m^2)'][pl_items.index("Payload")] = df['Mass (kg)'][pl_items.index("Payload")] * pow(df['Height (m)'][pl_items.index("Payload")],2)/12 
                    df['J0 (kg m^2)'][pl_items.index("PLF")] = df['Mass (kg)'][pl_items.index("PLF")] * (pow(step.r,2)/4 + pow(df['Height (m)'][pl_items.index("PLF")],2)/18)
                
                row += len(l_items)

            elif step.propulsion == 'Solid':
                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['J0 (kg m^2)'][row + s_items.index("Nose Cone")] = df['Mass (kg)'][row + s_items.index("Nose Cone")] * (pow(step.r,2)/4 + pow(df['Height (m)'][row + s_items.index("Nose Cone")],2)/18)    
                    df['J0 (kg m^2)'][row + s_items.index("Forward Skirt")] = df['Mass (kg)'][row + s_items.index("Forward Skirt")] * ( (pow(step.r,2) + pow(self.listOfSteps[i+1].r,2))/4 + pow(df['Height (m)'][row + s_items.index("Forward Skirt")],2)/18 * (1+(2*step.r*self.listOfSteps[i+1].r)/(pow(step.r + self.listOfSteps[i+1].r,2))))
                elif i == len(self.listOfSteps) - 1: # if last step                  
                    df['J0 (kg m^2)'][row + s_items.index("Forward Skirt")] = df['Mass (kg)'][row + s_items.index("Forward Skirt")] - self.CM_full
                    df['J0 (kg m^2)'][pl_items.index("PAF")] = df['Mass (kg)'][pl_items.index("PAF")] * pow(step.r,2)/2       
                    df['J0 (kg m^2)'][pl_items.index("Payload")] = df['Mass (kg)'][pl_items.index("Payload")] * pow(df['Height (m)'][pl_items.index("Payload")],2)/12             
                    df['J0 (kg m^2)'][pl_items.index("PLF")] = df['Mass (kg)'][pl_items.index("PLF")] * (pow(step.r,2)/4 + pow(df['Height (m)'][pl_items.index("PLF")],2)/18)               

                df['J0 (kg m^2)'][row + s_items.index("Avionics")] = df['Mass (kg)'][row + s_items.index("Avionics")] * pow(step.r,2)/2
                df['J0 (kg m^2)'][row + s_items.index("Wiring")] = df['Mass (kg)'][row + s_items.index("Wiring")] * pow(df['Height (m)'][row + s_items.index("Wiring")],2)/12                     
                df['J0 (kg m^2)'][row + s_items.index("Pressurant Tank")] = df['Mass (kg)'][row + s_items.index("Pressurant Tank")] * pow(step.r,2)
                df['J0 (kg m^2)'][row + s_items.index("SRM Dome Top")] = df['Mass (kg)'][row + s_items.index("SRM Dome Top")] * pow(step.r,2)
                df['J0 (kg m^2)'][row + s_items.index("Solid Propellant Casing")] = df['Mass (kg)'][row + s_items.index("Solid Propellant Casing")] * ( pow(step.r,2)/2 + pow(df['Height (m)'][row + s_items.index("Solid Propellant Casing")],2)/12 )
                df['J0 (kg m^2)'][row + s_items.index("SRM Dome Bottom")] = df['Mass (kg)'][row + s_items.index("SRM Dome Bottom")] * pow(step.r,2)
                df['J0 (kg m^2)'][row + s_items.index("Solid Propellant Residual")] = 0                 
                df['J0 (kg m^2)'][row + s_items.index("Aft Skirt")] = df['Mass (kg)'][row + s_items.index("Aft Skirt")] * ( pow(step.r,2) + pow(df['Height (m)'][row + s_items.index("Aft Skirt")],2)/12) 
                df['J0 (kg m^2)'][row + s_items.index("Gimballs")] = 0                 
                df['J0 (kg m^2)'][row + s_items.index("Nozzle")] = 0                  
                df['J0 (kg m^2)'][row + s_items.index("Solid Propellant")] = 0                 
                row += len(s_items)

    # init J0's - double check some values such as payload, fairing, propellant tanks etc
    def initJ0s(self, df, pl_items, l_items, s_items): 
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            
            if step.propulsion == "Liquid":
                df['J0 (kg m^2)'][row + l_items.index("Avionics")] = self.J0ThinRing(df['Mass (kg)'][row + l_items.index("Avionics")], step.r)
                df['J0 (kg m^2)'][row + l_items.index("Wiring")] = self.J0SolidRod(df['Mass (kg)'][row + l_items.index("Wiring")], df['Height (m)'][row + l_items.index("Wiring")])
                df['J0 (kg m^2)'][row + l_items.index("Fuel Dome Top")] = self.J0ThinEllipsoid(df['Mass (kg)'][row + l_items.index("Fuel Dome Top")], step.r, step.dome_f[0])
                df['J0 (kg m^2)'][row + l_items.index("Fuel Cylinder")] = self.J0ThinCyl(df['Mass (kg)'][row + l_items.index("Fuel Cylinder")], step.r, step.cyl_f[0])
                df['J0 (kg m^2)'][row + l_items.index("Fuel Dome Bottom")] = self.J0ThinEllipsoid(df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")], step.r, step.dome_f[0])
                SA_F_Tank = step.cyl_f[1] + 2 * step.dome_f[1]
                m_ins_f = df['Mass (kg)'][row + l_items.index("Fuel Insulation")]
                m_ins_f_dome = step.dome_f[1]/SA_F_Tank*m_ins_f
                m_ins_f_cyl = step.cyl_f[1]/SA_F_Tank*m_ins_f
                if m_ins_f == 0:
                    df['J0 (kg m^2)'][row + l_items.index("Fuel Insulation")] = 0
                else:
                    df['J0 (kg m^2)'][row + l_items.index("Fuel Insulation")] = (2*m_ins_f_dome*self.J0ThinEllipsoid(m_ins_f_dome, step.r, step.dome_f[0]) + m_ins_f_cyl*self.J0ThinCyl(m_ins_f_cyl, step.r, step.cyl_f[0]))/m_ins_f
                df['J0 (kg m^2)'][row + l_items.index("Fuel Residual")] = 0   
                df['J0 (kg m^2)'][row + l_items.index("Intertank")] = self.J0ThinCyl(df['Mass (kg)'][row + l_items.index("Intertank")], step.r, step.intertank[0])
                df['J0 (kg m^2)'][row + l_items.index("Ox Dome Top")] = self.J0ThinEllipsoid(df['Mass (kg)'][row + l_items.index("Ox Dome Top")], step.r, step.dome_ox[0])
                df['J0 (kg m^2)'][row + l_items.index("Ox Cylinder")] = self.J0ThinCyl(df['Mass (kg)'][row + l_items.index("Ox Cylinder")], step.r, step.cyl_ox[0])
                df['J0 (kg m^2)'][row + l_items.index("Ox Dome Bottom")] = self.J0ThinEllipsoid(df['Mass (kg)'][row + l_items.index("Ox Dome Bottom")], step.r, step.dome_ox[0])
                SA_Ox_Tank = step.cyl_ox[1] + 2 * step.dome_ox[1]
                m_ins_ox = df['Mass (kg)'][row + l_items.index("Ox Insulation")]
                m_ins_ox_dome = step.dome_ox[1]/SA_Ox_Tank*m_ins_ox
                m_ins_ox_cyl = step.cyl_ox[1]/SA_Ox_Tank*m_ins_ox
                df['J0 (kg m^2)'][row + l_items.index("Ox Insulation")] = (2*m_ins_ox_dome*self.J0ThinEllipsoid(m_ins_ox_dome, step.r, step.dome_ox[0]) + m_ins_ox_cyl*self.J0ThinCyl(m_ins_ox_cyl, step.r, step.cyl_ox[0]))/m_ins_ox
                df['J0 (kg m^2)'][row + l_items.index("Ox Residual")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Pressurant Tank")] = self.J0ThinHemisphere(df['Mass (kg)'][row + l_items.index("Pressurant Tank")], step.r)
                df['J0 (kg m^2)'][row + l_items.index("Aft Skirt")] = self.J0ThinCyl(df['Mass (kg)'][row + l_items.index("Aft Skirt")], step.r, step.aft_skirt[0])
                df['J0 (kg m^2)'][row + l_items.index("Thrust Structure")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Gimballs")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Engines")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Fuel")] = 0
                df['J0 (kg m^2)'][row + l_items.index("Oxidizer")] = 0

                if i < len(self.listOfSteps) - 1: # if not last step
                    # Interstage ( using formual for centroid of trapezoid)
                    df['J0 (kg m^2)'][row + l_items.index("Forward Skirt")] = self.J0ThinTrap(df['Mass (kg)'][row + l_items.index("Forward Skirt")], step.r, self.listOfSteps[i+1].r, df['Height (m)'][row + l_items.index("Forward Skirt")])
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['J0 (kg m^2)'][row + l_items.index("Forward Skirt")] = self.J0ThinCyl(df['Mass (kg)'][row + l_items.index("Forward Skirt")], step.r, step.fwd_skirt[0])
                    df['J0 (kg m^2)'][pl_items.index("PAF")] = self.J0ThinRing(df['Mass (kg)'][pl_items.index("PAF")], step.r)
                    df['J0 (kg m^2)'][pl_items.index("Payload")] = self.J0SolidCyl(df['Mass (kg)'][pl_items.index("Payload")], step.r, step.fairing[0][0])
                    last_step = self.listOfSteps[len(self.listOfSteps)-1]
                    SA_PLF = sum(last_step.fairing[1])
                    m_PLF = df['Mass (kg)'][pl_items.index("PLF")]
                    m_PLF_cyl = step.fairing[1][0]/SA_PLF*m_PLF
                    m_PLF_cone = step.fairing[1][1]/SA_PLF*m_PLF
                    df['J0 (kg m^2)'][pl_items.index("PLF")] = (m_PLF_cyl*self.J0ThinCyl(m_PLF_cyl, step.r, step.fairing[0][0]) + m_PLF_cone*self.J0ThinCone(m_PLF_cone, step.r, step.fairing[0][1]))/m_PLF
                
                row += len(l_items)

            elif step.propulsion == 'Solid':
                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['J0 (kg m^2)'][row + s_items.index("Nose Cone")] = self.J0ThinCone(df['Mass (kg)'][row + s_items.index("Nose Cone")], step.r, step.fairing[0])
                    df['J0 (kg m^2)'][row + s_items.index("Forward Skirt")] = self.J0ThinTrap(df['Mass (kg)'][row + s_items.index("Forward Skirt")], step.r, self.listOfSteps[i+1].r, df['Height (m)'][row + s_items.index("Forward Skirt")])
                elif i == len(self.listOfSteps) - 1: # if last step                  
                    df['J0 (kg m^2)'][row + s_items.index("Forward Skirt")] = self.J0ThinCyl(df['Mass (kg)'][row + s_items.index("Forward Skirt")], step.r, step.fwd_skirt[0])
                    df['J0 (kg m^2)'][pl_items.index("PAF")] = self.J0ThinRing(df['Mass (kg)'][pl_items.index("PAF")], step.r)
                    df['J0 (kg m^2)'][pl_items.index("Payload")] = self.J0SolidCyl(df['Mass (kg)'][pl_items.index("Payload")], step.r, step.fairing[0][0])
                    m_PLF = df['Mass (kg)'][pl_items.index("PLF")]
                    m_PLF_cyl = step.fairing[1][0]/SA_PLF*m_PLF
                    m_PLF_cone = step.fairing[1][1]/SA_PLF*m_PLF
                    df['J0 (kg m^2)'][pl_items.index("PLF")] = (m_PLF_cyl*self.J0ThinCyl(m_PLF_cyl, step.r, step.fairing[0][0]) + m_PLF_cone*self.J0ThinCone(m_PLF_cone, step.r, step.fairing[0][1]))/m_PLF

                df['J0 (kg m^2)'][row + s_items.index("Avionics")] = self.J0ThinRing(df['Mass (kg)'][row + s_items.index("Avionics")], step.r)
                df['J0 (kg m^2)'][row + s_items.index("Wiring")] = self.J0SolidRod(df['Mass (kg)'][row + s_items.index("Wiring")], df['Height (m)'][row + s_items.index("Wiring")])
                df['J0 (kg m^2)'][row + s_items.index("Pressurant Tank")] = 2*self.J0ThinHemisphere(df['Mass (kg)'][row + s_items.index("Pressurant Tank")], step.r)
                df['J0 (kg m^2)'][row + s_items.index("SRM Dome Top")] = self.J0ThinEllipsoid(df['Mass (kg)'][row + s_items.index("SRM Dome Top")], step.r, step.dome_f[0])
                df['J0 (kg m^2)'][row + s_items.index("Solid Propellant Casing")] = self.J0ThinCyl(df['Mass (kg)'][row + s_items.index("Solid Propellant Casing")], step.r, step.srm_casing[0])
                df['J0 (kg m^2)'][row + s_items.index("SRM Dome Bottom")] = self.J0ThinEllipsoid(df['Mass (kg)'][row + s_items.index("SRM Dome Bottom")], step.r, step.dome_f[0])
                df['J0 (kg m^2)'][row + s_items.index("Solid Propellant Residual")] = 0
                df['J0 (kg m^2)'][row + s_items.index("Aft Skirt")] = self.J0ThinCyl(df['Mass (kg)'][row + s_items.index("Aft Skirt")], step.r, step.aft_skirt[0])
                df['J0 (kg m^2)'][row + s_items.index("Gimballs")] = 0
                df['J0 (kg m^2)'][row + s_items.index("Nozzle")] = 0
                SA_SRM = step.dome_f[1]*2 + step.srm_casing[1]
                m_solid_prop = df['Mass (kg)'][row + s_items.index("Solid Propellant")]
                m_dome = step.dome_f[1]/SA_SRM*m_solid_prop
                m_casing = step.srm_casing[1]/SA_SRM*m_solid_prop
                # J0_solid propellant = (2*m_dome*J0_dome + m_casing*J0_casing)/m_total
                df['J0 (kg m^2)'][row + s_items.index("Solid Propellant")] = (2*m_dome*self.J0SolidEllipsoid(m_dome, step.r, step.dome_f[0]) + m_casing*self.J0SolidCyl(m_casing, step.r, step.srm_casing[0]))/m_solid_prop
                row += len(s_items)

    def initmCMs(self, df, pl_items, l_items, s_items):
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            
            if step.propulsion == "Liquid":
                df['m*CM^2 (kg m^2)'][row + l_items.index("Avionics")] = df['Mass (kg)'][row + l_items.index("Avionics")] * pow(df['Distance from CM (m)'][row + l_items.index("Avionics")],2)     
                df['m*CM^2 (kg m^2)'][row + l_items.index("Wiring")] = df['Mass (kg)'][row + l_items.index("Wiring")] * pow(df['Distance from CM (m)'][row + l_items.index("Wiring")],2)                
                df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Dome Top")] = df['Mass (kg)'][row + l_items.index("Fuel Dome Top")] * pow(df['Distance from CM (m)'][row + l_items.index("Fuel Dome Top")],2)     
                df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Cylinder")] = df['Mass (kg)'][row + l_items.index("Fuel Cylinder")] * pow(df['Distance from CM (m)'][row + l_items.index("Fuel Cylinder")],2)      
                df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Dome Bottom")] = df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")] * pow(df['Distance from CM (m)'][row + l_items.index("Fuel Dome Bottom")],2)    
                df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Insulation")] = df['Mass (kg)'][row + l_items.index("Fuel Insulation")] * pow(df['Distance from CM (m)'][row + l_items.index("Fuel Insulation")],2)      
                df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Residual")] = df['Mass (kg)'][row + l_items.index("Fuel Residual")] * pow(df['Distance from CM (m)'][row + l_items.index("Fuel Residual")],2)                 
                df['m*CM^2 (kg m^2)'][row + l_items.index("Intertank")] = df['Mass (kg)'][row + l_items.index("Intertank")] * pow(df['Distance from CM (m)'][row + l_items.index("Intertank")],2)     
                df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Dome Top")] = df['Mass (kg)'][row + l_items.index("Ox Dome Top")] * pow(df['Distance from CM (m)'][row + l_items.index("Ox Dome Top")],2)     
                df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Cylinder")] = df['Mass (kg)'][row + l_items.index("Ox Cylinder")] * pow(df['Distance from CM (m)'][row + l_items.index("Ox Cylinder")],2)
                df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Dome Bottom")] = df['Mass (kg)'][row + l_items.index("Ox Dome Bottom")] * pow(df['Distance from CM (m)'][row + l_items.index("Ox Dome Bottom")],2)   
                df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Insulation")] = df['Mass (kg)'][row + l_items.index("Ox Insulation")] * pow(df['Distance from CM (m)'][row + l_items.index("Ox Insulation")],2)   
                df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Residual")] = df['Mass (kg)'][row + l_items.index("Ox Residual")] * pow(df['Distance from CM (m)'][row + l_items.index("Ox Residual")],2)                
                df['m*CM^2 (kg m^2)'][row + l_items.index("Pressurant Tank")] = df['Mass (kg)'][row + l_items.index("Pressurant Tank")] * pow(df['Distance from CM (m)'][row + l_items.index("Pressurant Tank")],2)
                df['m*CM^2 (kg m^2)'][row + l_items.index("Aft Skirt")] = df['Mass (kg)'][row + l_items.index("Aft Skirt")] * pow(df['Distance from CM (m)'][row + l_items.index("Aft Skirt")],2) 
                df['m*CM^2 (kg m^2)'][row + l_items.index("Thrust Structure")] = df['Mass (kg)'][row + l_items.index("Thrust Structure")] * pow(df['Distance from CM (m)'][row + l_items.index("Thrust Structure")],2)
                df['m*CM^2 (kg m^2)'][row + l_items.index("Gimballs")] = df['Mass (kg)'][row + l_items.index("Gimballs")] * pow(df['Distance from CM (m)'][row + l_items.index("Gimballs")],2)
                df['m*CM^2 (kg m^2)'][row + l_items.index("Engines")] = df['Mass (kg)'][row + l_items.index("Engines")] * pow(df['Distance from CM (m)'][row + l_items.index("Engines")],2)
                df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel")] = df['Mass (kg)'][row + l_items.index("Fuel")] * pow(df['Distance from CM (m)'][row + l_items.index("Fuel")],2)
                df['m*CM^2 (kg m^2)'][row + l_items.index("Oxidizer")] = df['Mass (kg)'][row + l_items.index("Oxidizer")] * pow(df['Distance from CM (m)'][row + l_items.index("Oxidizer")],2)

                if i < len(self.listOfSteps) - 1: # if not last step
                   # Interstage ( using formual for centroid of trapezoid)
                    df['m*CM^2 (kg m^2)'][row + l_items.index("Forward Skirt")] = df['Mass (kg)'][row + l_items.index("Forward Skirt")] * pow(df['Distance from CM (m)'][row + l_items.index("Forward Skirt")],2)
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['m*CM^2 (kg m^2)'][row + l_items.index("Forward Skirt")] = df['Mass (kg)'][row + l_items.index("Forward Skirt")] * pow(df['Distance from CM (m)'][row + l_items.index("Forward Skirt")],2)          
                    df['m*CM^2 (kg m^2)'][pl_items.index("PAF")] = df['Mass (kg)'][pl_items.index("PAF")] * pow(df['Distance from CM (m)'][pl_items.index("PAF")],2)                  
                    df['m*CM^2 (kg m^2)'][pl_items.index("Payload")] = df['Mass (kg)'][pl_items.index("Payload")] * pow(df['Distance from CM (m)'][pl_items.index("Payload")],2) 
                    df['m*CM^2 (kg m^2)'][pl_items.index("PLF")] = df['Mass (kg)'][pl_items.index("PLF")] * pow(df['Distance from CM (m)'][pl_items.index("PLF")],2)
                
                row += len(l_items)

            elif step.propulsion == 'Solid':

                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['m*CM^2 (kg m^2)'][row + s_items.index("Nose Cone")] = df['Mass (kg)'][row + s_items.index("Nose Cone")] * pow(df['Distance from CM (m)'][row + s_items.index("Nose Cone")],2)         
                    df['m*CM^2 (kg m^2)'][row + s_items.index("Forward Skirt")] = df['Mass (kg)'][row + s_items.index("Forward Skirt")] * pow(df['Distance from CM (m)'][row + s_items.index("Forward Skirt")],2)
                elif i == len(self.listOfSteps) - 1: # if last step        
                    df['m*CM^2 (kg m^2)'][row + s_items.index("Forward Skirt")] = df['Mass (kg)'][row + s_items.index("Forward Skirt")] * pow(df['Distance from CM (m)'][row + s_items.index("Forward Skirt")],2)
                    df['m*CM^2 (kg m^2)'][pl_items.index("PAF")] = df['Mass (kg)'][pl_items.index("PAF")] * pow(df['Distance from CM (m)'][pl_items.index("PAF")],2)                   
                    df['m*CM^2 (kg m^2)'][pl_items.index("Payload")] = df['Mass (kg)'][pl_items.index("Payload")] * pow(df['Distance from CM (m)'][pl_items.index("Payload")],2)     
                    df['m*CM^2 (kg m^2)'][pl_items.index("PLF")] = df['Mass (kg)'][pl_items.index("PLF")] * pow(df['Distance from CM (m)'][pl_items.index("PLF")],2)    

                df['m*CM^2 (kg m^2)'][row + s_items.index("Avionics")] = df['Mass (kg)'][row + s_items.index("Avionics")] * pow(df['Distance from CM (m)'][row + s_items.index("Avionics")],2)
                df['m*CM^2 (kg m^2)'][row + s_items.index("Wiring")] = df['Mass (kg)'][row + s_items.index("Wiring")] * pow(df['Distance from CM (m)'][row + s_items.index("Wiring")],2)                     
                df['m*CM^2 (kg m^2)'][row + s_items.index("Pressurant Tank")] = df['Mass (kg)'][row + s_items.index("Pressurant Tank")] * pow(df['Distance from CM (m)'][row + s_items.index("Pressurant Tank")],2)
                df['m*CM^2 (kg m^2)'][row + s_items.index("SRM Dome Top")] = df['Mass (kg)'][row + s_items.index("SRM Dome Top")] * pow(df['Distance from CM (m)'][row + s_items.index("SRM Dome Top")],2) 
                df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant Casing")] = df['Mass (kg)'][row + s_items.index("Solid Propellant Casing")] * pow(df['Distance from CM (m)'][row + s_items.index("Solid Propellant Casing")],2)
                df['m*CM^2 (kg m^2)'][row + s_items.index("SRM Dome Bottom")] = df['Mass (kg)'][row + s_items.index("SRM Dome Bottom")] * pow(df['Distance from CM (m)'][row + s_items.index("SRM Dome Bottom")],2) 
                df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant Residual")] = df['Mass (kg)'][row + s_items.index("Solid Propellant Residual")] * pow(df['Distance from CM (m)'][row + s_items.index("Solid Propellant Residual")],2)                 
                df['m*CM^2 (kg m^2)'][row + s_items.index("Aft Skirt")] = df['Mass (kg)'][row + s_items.index("Aft Skirt")] * pow(df['Distance from CM (m)'][row + s_items.index("Aft Skirt")],2) 
                df['m*CM^2 (kg m^2)'][row + s_items.index("Gimballs")] = df['Mass (kg)'][row + s_items.index("Gimballs")] * pow(df['Distance from CM (m)'][row + s_items.index("Gimballs")],2)                 
                df['m*CM^2 (kg m^2)'][row + s_items.index("Nozzle")] = df['Mass (kg)'][row + s_items.index("Nozzle")] * pow(df['Distance from CM (m)'][row + s_items.index("Nozzle")],2)                  
                df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant")] = df['Mass (kg)'][row + s_items.index("Solid Propellant")] * pow(df['Distance from CM (m)'][row + s_items.index("Solid Propellant")],2)                 
                
                row += len(s_items)

    def initJPitchYawOld(self, df, pl_items, l_items, s_items):
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            if step.propulsion == "Liquid":
                df['Jpitch/yaw'][row + l_items.index("Avionics")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Avionics")] + df['Distance from CM (m)'][row + l_items.index("Avionics")]     
                df['Jpitch/yaw'][row + l_items.index("Wiring")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Wiring")] + df['Distance from CM (m)'][row + l_items.index("Wiring")]                
                df['Jpitch/yaw'][row + l_items.index("Fuel Dome Top")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Dome Top")] + df['Distance from CM (m)'][row + l_items.index("Fuel Dome Top")]     
                df['Jpitch/yaw'][row + l_items.index("Fuel Cylinder")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Cylinder")] + df['Distance from CM (m)'][row + l_items.index("Fuel Cylinder")]      
                df['Jpitch/yaw'][row + l_items.index("Fuel Dome Bottom")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Dome Bottom")] + df['Distance from CM (m)'][row + l_items.index("Fuel Dome Bottom")]    
                df['Jpitch/yaw'][row + l_items.index("Fuel Insulation")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Insulation")] + df['Distance from CM (m)'][row + l_items.index("Fuel Insulation")]      
                df['Jpitch/yaw'][row + l_items.index("Fuel Residual")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Residual")] + df['Distance from CM (m)'][row + l_items.index("Fuel Residual")]                 
                df['Jpitch/yaw'][row + l_items.index("Intertank")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Intertank")] + df['Distance from CM (m)'][row + l_items.index("Intertank")]     
                df['Jpitch/yaw'][row + l_items.index("Ox Dome Top")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Dome Top")] + df['Distance from CM (m)'][row + l_items.index("Ox Dome Top")]     
                df['Jpitch/yaw'][row + l_items.index("Ox Cylinder")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Cylinder")] + df['Distance from CM (m)'][row + l_items.index("Ox Cylinder")]
                df['Jpitch/yaw'][row + l_items.index("Ox Dome Bottom")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Dome Bottom")] + df['Distance from CM (m)'][row + l_items.index("Ox Dome Bottom")]   
                df['Jpitch/yaw'][row + l_items.index("Ox Insulation")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Insulation")] + df['Distance from CM (m)'][row + l_items.index("Ox Insulation")]   
                df['Jpitch/yaw'][row + l_items.index("Ox Residual")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Residual")] + df['Distance from CM (m)'][row + l_items.index("Ox Residual")]                
                df['Jpitch/yaw'][row + l_items.index("Pressurant Tank")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Pressurant Tank")] + df['Distance from CM (m)'][row + l_items.index("Pressurant Tank")]
                df['Jpitch/yaw'][row + l_items.index("Aft Skirt")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Aft Skirt")] + df['Distance from CM (m)'][row + l_items.index("Aft Skirt")] 
                df['Jpitch/yaw'][row + l_items.index("Thrust Structure")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Thrust Structure")] + df['Distance from CM (m)'][row + l_items.index("Thrust Structure")]
                df['Jpitch/yaw'][row + l_items.index("Gimballs")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Gimballs")] + df['Distance from CM (m)'][row + l_items.index("Gimballs")]
                df['Jpitch/yaw'][row + l_items.index("Engines")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Engines")] + df['Distance from CM (m)'][row + l_items.index("Engines")]
                df['Jpitch/yaw'][row + l_items.index("Fuel")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel")] + df['Distance from CM (m)'][row + l_items.index("Fuel")]
                df['Jpitch/yaw'][row + l_items.index("Oxidizer")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Oxidizer")] + df['Distance from CM (m)'][row + l_items.index("Oxidizer")]

                if i < len(self.listOfSteps) - 1: # if not last step
                   # Interstage ( using formual for centroid of trapezoid)
                    df['Jpitch/yaw'][row + l_items.index("Forward Skirt")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Forward Skirt")] + df['Distance from CM (m)'][row + l_items.index("Forward Skirt")]
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['Jpitch/yaw'][row + l_items.index("Forward Skirt")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Forward Skirt")] + df['Distance from CM (m)'][row + l_items.index("Forward Skirt")]       
                    df['Jpitch/yaw'][pl_items.index("PAF")] = df['m*CM^2 (kg m^2)'][pl_items.index("PAF")] + df['Distance from CM (m)'][pl_items.index("PAF")]              
                    df['Jpitch/yaw'][pl_items.index("Payload")] = df['m*CM^2 (kg m^2)'][pl_items.index("Payload")] + df['Distance from CM (m)'][pl_items.index("Payload")]              
                    df['Jpitch/yaw'][pl_items.index("PLF")] = df['m*CM^2 (kg m^2)'][pl_items.index("PLF")] + df['Distance from CM (m)'][pl_items.index("PLF")]             
                
                row += len(l_items)

            elif step.propulsion == 'Solid':
                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['Jpitch/yaw'][row + s_items.index("Nose Cone")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Nose Cone")] + df['Distance from CM (m)'][row + s_items.index("Nose Cone")]
                    df['Jpitch/yaw'][row + s_items.index("Forward Skirt")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Forward Skirt")] + df['Distance from CM (m)'][row + s_items.index("Forward Skirt")]
                elif i == len(self.listOfSteps) - 1: # if last step      
                    df['Jpitch/yaw'][row + s_items.index("Forward Skirt")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Forward Skirt")] + df['Distance from CM (m)'][row + s_items.index("Forward Skirt")]
                    df['Jpitch/yaw'][pl_items.index("PAF")] = df['m*CM^2 (kg m^2)'][pl_items.index("PAF")] + df['Distance from CM (m)'][pl_items.index("PAF")]                  
                    df['Jpitch/yaw'][pl_items.index("Payload")] = df['m*CM^2 (kg m^2)'][pl_items.index("Payload")] + df['Distance from CM (m)'][pl_items.index("Payload")]                  
                    df['Jpitch/yaw'][pl_items.index("PLF")] = df['m*CM^2 (kg m^2)'][pl_items.index("PLF")] + df['Distance from CM (m)'][pl_items.index("PLF")]                 

                df['Jpitch/yaw'][row + s_items.index("Avionics")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Avionics")] + df['Distance from CM (m)'][row + s_items.index("Avionics")]
                df['Jpitch/yaw'][row + s_items.index("Wiring")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Wiring")] + df['Distance from CM (m)'][row + s_items.index("Wiring")]
                df['Jpitch/yaw'][row + s_items.index("Pressurant Tank")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Pressurant Tank")] + df['Distance from CM (m)'][row + s_items.index("Pressurant Tank")]
                df['Jpitch/yaw'][row + s_items.index("SRM Dome Top")] = df['m*CM^2 (kg m^2)'][row + s_items.index("SRM Dome Top")] + df['Distance from CM (m)'][row + s_items.index("SRM Dome Top")] 
                df['Jpitch/yaw'][row + s_items.index("Solid Propellant Casing")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant Casing")] + df['Distance from CM (m)'][row + s_items.index("Solid Propellant Casing")] 
                df['Jpitch/yaw'][row + s_items.index("SRM Dome Bottom")] = df['m*CM^2 (kg m^2)'][row + s_items.index("SRM Dome Bottom")] + df['Distance from CM (m)'][row + s_items.index("SRM Dome Bottom")] 
                df['Jpitch/yaw'][row + s_items.index("Solid Propellant Residual")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant Residual")] + df['Distance from CM (m)'][row + s_items.index("Solid Propellant Residual")]                 
                df['Jpitch/yaw'][row + s_items.index("Aft Skirt")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Aft Skirt")] + df['Distance from CM (m)'][row + s_items.index("Aft Skirt")] 
                df['Jpitch/yaw'][row + s_items.index("Gimballs")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Gimballs")] + df['Distance from CM (m)'][row + s_items.index("Gimballs")]                 
                df['Jpitch/yaw'][row + s_items.index("Nozzle")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Nozzle")] + df['Distance from CM (m)'][row + s_items.index("Nozzle")]                  
                df['Jpitch/yaw'][row + s_items.index("Solid Propellant")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant")] + df['Distance from CM (m)'][row + s_items.index("Solid Propellant")]                 
                
                row += len(s_items)

    def initJPitchYaw(self, df, pl_items, l_items, s_items):
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            if step.propulsion == "Liquid":
                df['Jpitch/yaw'][row + l_items.index("Avionics")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Avionics")] + df['J0 (kg m^2)'][row + l_items.index("Avionics")]     
                df['Jpitch/yaw'][row + l_items.index("Wiring")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Wiring")] + df['J0 (kg m^2)'][row + l_items.index("Wiring")]                
                df['Jpitch/yaw'][row + l_items.index("Fuel Dome Top")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Dome Top")] + df['J0 (kg m^2)'][row + l_items.index("Fuel Dome Top")]     
                df['Jpitch/yaw'][row + l_items.index("Fuel Cylinder")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Cylinder")] + df['J0 (kg m^2)'][row + l_items.index("Fuel Cylinder")]      
                df['Jpitch/yaw'][row + l_items.index("Fuel Dome Bottom")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Dome Bottom")] + df['J0 (kg m^2)'][row + l_items.index("Fuel Dome Bottom")]    
                df['Jpitch/yaw'][row + l_items.index("Fuel Insulation")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Insulation")] + df['J0 (kg m^2)'][row + l_items.index("Fuel Insulation")]      
                df['Jpitch/yaw'][row + l_items.index("Fuel Residual")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel Residual")] + df['J0 (kg m^2)'][row + l_items.index("Fuel Residual")]                 
                df['Jpitch/yaw'][row + l_items.index("Intertank")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Intertank")] + df['J0 (kg m^2)'][row + l_items.index("Intertank")]     
                df['Jpitch/yaw'][row + l_items.index("Ox Dome Top")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Dome Top")] + df['J0 (kg m^2)'][row + l_items.index("Ox Dome Top")]     
                df['Jpitch/yaw'][row + l_items.index("Ox Cylinder")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Cylinder")] + df['J0 (kg m^2)'][row + l_items.index("Ox Cylinder")]
                df['Jpitch/yaw'][row + l_items.index("Ox Dome Bottom")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Dome Bottom")] + df['J0 (kg m^2)'][row + l_items.index("Ox Dome Bottom")]   
                df['Jpitch/yaw'][row + l_items.index("Ox Insulation")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Insulation")] + df['J0 (kg m^2)'][row + l_items.index("Ox Insulation")]   
                df['Jpitch/yaw'][row + l_items.index("Ox Residual")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Ox Residual")] + df['J0 (kg m^2)'][row + l_items.index("Ox Residual")]                
                df['Jpitch/yaw'][row + l_items.index("Pressurant Tank")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Pressurant Tank")] + df['J0 (kg m^2)'][row + l_items.index("Pressurant Tank")]
                df['Jpitch/yaw'][row + l_items.index("Aft Skirt")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Aft Skirt")] + df['J0 (kg m^2)'][row + l_items.index("Aft Skirt")] 
                df['Jpitch/yaw'][row + l_items.index("Thrust Structure")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Thrust Structure")] + df['J0 (kg m^2)'][row + l_items.index("Thrust Structure")]
                df['Jpitch/yaw'][row + l_items.index("Gimballs")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Gimballs")] + df['J0 (kg m^2)'][row + l_items.index("Gimballs")]
                df['Jpitch/yaw'][row + l_items.index("Engines")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Engines")] + df['J0 (kg m^2)'][row + l_items.index("Engines")]
                df['Jpitch/yaw'][row + l_items.index("Fuel")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Fuel")] + df['J0 (kg m^2)'][row + l_items.index("Fuel")]
                df['Jpitch/yaw'][row + l_items.index("Oxidizer")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Oxidizer")] + df['J0 (kg m^2)'][row + l_items.index("Oxidizer")]

                if i < len(self.listOfSteps) - 1: # if not last step
                   # Interstage ( using formual for centroid of trapezoid)
                    df['Jpitch/yaw'][row + l_items.index("Forward Skirt")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Forward Skirt")] + df['J0 (kg m^2)'][row + l_items.index("Forward Skirt")]
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['Jpitch/yaw'][row + l_items.index("Forward Skirt")] = df['m*CM^2 (kg m^2)'][row + l_items.index("Forward Skirt")] + df['J0 (kg m^2)'][row + l_items.index("Forward Skirt")]       
                    df['Jpitch/yaw'][pl_items.index("PAF")] = df['m*CM^2 (kg m^2)'][pl_items.index("PAF")] + df['J0 (kg m^2)'][pl_items.index("PAF")]              
                    df['Jpitch/yaw'][pl_items.index("Payload")] = df['m*CM^2 (kg m^2)'][pl_items.index("Payload")] + df['J0 (kg m^2)'][pl_items.index("Payload")]              
                    df['Jpitch/yaw'][pl_items.index("PLF")] = df['m*CM^2 (kg m^2)'][pl_items.index("PLF")] + df['J0 (kg m^2)'][pl_items.index("PLF")]             
                
                row += len(l_items)

            elif step.propulsion == 'Solid':
                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['Jpitch/yaw'][row + s_items.index("Nose Cone")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Nose Cone")] + df['J0 (kg m^2)'][row + s_items.index("Nose Cone")]
                    df['Jpitch/yaw'][row + s_items.index("Forward Skirt")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Forward Skirt")] + df['J0 (kg m^2)'][row + s_items.index("Forward Skirt")]
                elif i == len(self.listOfSteps) - 1: # if last step      
                    df['Jpitch/yaw'][row + s_items.index("Forward Skirt")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Forward Skirt")] + df['J0 (kg m^2)'][row + s_items.index("Forward Skirt")]
                    df['Jpitch/yaw'][pl_items.index("PAF")] = df['m*CM^2 (kg m^2)'][pl_items.index("PAF")] + df['J0 (kg m^2)'][pl_items.index("PAF")]                  
                    df['Jpitch/yaw'][pl_items.index("Payload")] = df['m*CM^2 (kg m^2)'][pl_items.index("Payload")] + df['J0 (kg m^2)'][pl_items.index("Payload")]                  
                    df['Jpitch/yaw'][pl_items.index("PLF")] = df['m*CM^2 (kg m^2)'][pl_items.index("PLF")] + df['J0 (kg m^2)'][pl_items.index("PLF")]                 

                df['Jpitch/yaw'][row + s_items.index("Avionics")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Avionics")] + df['J0 (kg m^2)'][row + s_items.index("Avionics")]
                df['Jpitch/yaw'][row + s_items.index("Wiring")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Wiring")] + df['J0 (kg m^2)'][row + s_items.index("Wiring")]
                df['Jpitch/yaw'][row + s_items.index("Pressurant Tank")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Pressurant Tank")] + df['J0 (kg m^2)'][row + s_items.index("Pressurant Tank")]
                df['Jpitch/yaw'][row + s_items.index("SRM Dome Top")] = df['m*CM^2 (kg m^2)'][row + s_items.index("SRM Dome Top")] + df['J0 (kg m^2)'][row + s_items.index("SRM Dome Top")] 
                df['Jpitch/yaw'][row + s_items.index("Solid Propellant Casing")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant Casing")] + df['J0 (kg m^2)'][row + s_items.index("Solid Propellant Casing")] 
                df['Jpitch/yaw'][row + s_items.index("SRM Dome Bottom")] = df['m*CM^2 (kg m^2)'][row + s_items.index("SRM Dome Bottom")] + df['J0 (kg m^2)'][row + s_items.index("SRM Dome Bottom")] 
                df['Jpitch/yaw'][row + s_items.index("Solid Propellant Residual")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant Residual")] + df['J0 (kg m^2)'][row + s_items.index("Solid Propellant Residual")]                 
                df['Jpitch/yaw'][row + s_items.index("Aft Skirt")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Aft Skirt")] + df['J0 (kg m^2)'][row + s_items.index("Aft Skirt")] 
                df['Jpitch/yaw'][row + s_items.index("Gimballs")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Gimballs")] + df['J0 (kg m^2)'][row + s_items.index("Gimballs")]                 
                df['Jpitch/yaw'][row + s_items.index("Nozzle")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Nozzle")] + df['J0 (kg m^2)'][row + s_items.index("Nozzle")]                  
                df['Jpitch/yaw'][row + s_items.index("Solid Propellant")] = df['m*CM^2 (kg m^2)'][row + s_items.index("Solid Propellant")] + df['J0 (kg m^2)'][row + s_items.index("Solid Propellant")]                 
                
                row += len(s_items)

    def initJRollOld(self, df, pl_items, l_items, s_items):
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            
            if step.propulsion == "Liquid":
                df['Jroll'][row + l_items.index("Avionics")] = df['Mass (kg)'][row + l_items.index("Avionics")] * pow(step.r,2)       
                df['Jroll'][row + l_items.index("Wiring")] = df['Mass (kg)'][row + l_items.index("Wiring")] * pow(step.r,2)     
                df['Jroll'][row + l_items.index("Fuel Dome Top")] = df['Mass (kg)'][row + l_items.index("Fuel Dome Top")] * pow(step.r,2)     
                df['Jroll'][row + l_items.index("Fuel Cylinder")] = df['Mass (kg)'][row + l_items.index("Fuel Cylinder")] * pow(step.r,2)      
                df['Jroll'][row + l_items.index("Fuel Dome Bottom")] = df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")] * pow(step.r,2)    
                df['Jroll'][row + l_items.index("Fuel Insulation")] = df['Mass (kg)'][row + l_items.index("Fuel Insulation")] * pow(step.r,2)      
                df['Jroll'][row + l_items.index("Fuel Residual")] = 0   
                df['Jroll'][row + l_items.index("Intertank")] = df['Mass (kg)'][row + l_items.index("Intertank")] * pow(step.r,2)     
                df['Jroll'][row + l_items.index("Ox Dome Top")] = df['Mass (kg)'][row + l_items.index("Ox Dome Top")] * pow(step.r,2)     
                df['Jroll'][row + l_items.index("Ox Cylinder")] = df['Mass (kg)'][row + l_items.index("Ox Cylinder")] * pow(step.r,2)
                df['Jroll'][row + l_items.index("Ox Dome Bottom")] = df['Mass (kg)'][row + l_items.index("Ox Dome Bottom")] * pow(step.r,2)   
                df['Jroll'][row + l_items.index("Ox Insulation")] = df['Mass (kg)'][row + l_items.index("Ox Insulation")] * pow(step.r,2)
                df['Jroll'][row + l_items.index("Ox Residual")] = 0  
                df['Jroll'][row + l_items.index("Pressurant Tank")] = df['Mass (kg)'][row + l_items.index("Pressurant Tank")] * pow(step.r,2)
                df['Jroll'][row + l_items.index("Ox Dome Bottom")] = df['Mass (kg)'][row + l_items.index("Ox Dome Bottom")] * pow(step.r,2)
                df['Jroll'][row + l_items.index("Aft Skirt")] = df['Mass (kg)'][row + l_items.index("Aft Skirt")] * pow(step.r,2) 
                df['Jroll'][row + l_items.index("Thrust Structure")] = df['Mass (kg)'][row + l_items.index("Thrust Structure")] * pow(step.r,2) 
                df['Jroll'][row + l_items.index("Gimballs")] = 0 
                df['Jroll'][row + l_items.index("Engines")] = 0 
                df['Jroll'][row + l_items.index("Fuel")] = 0 
                df['Jroll'][row + l_items.index("Oxidizer")] = 0 

                if i < len(self.listOfSteps) - 1: # if not last step
                   # Interstage ( using formual for centroid of trapezoid)
                    df['Jroll'][row + l_items.index("Forward Skirt")] = df['Mass (kg)'][row + l_items.index("Forward Skirt")] * pow(step.r,2)
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['Jroll'][row + l_items.index("Forward Skirt")] = df['Mass (kg)'][row + l_items.index("Forward Skirt")] * pow(step.r,2)          
                    df['Jroll'][pl_items.index("PAF")] = df['Mass (kg)'][pl_items.index("PAF")] * pow(step.r,2)                  
                    df['Jroll'][pl_items.index("Payload")] = df['Mass (kg)'][pl_items.index("Payload")] * pow(step.r,2) 
                    df['Jroll'][pl_items.index("PLF")] = df['Mass (kg)'][pl_items.index("PLF")] * pow(step.r,2)
                
                row += len(l_items)

            elif step.propulsion == 'Solid':
                df['Jroll'][row + s_items.index("Avionics")] = df['Mass (kg)'][row + s_items.index("Avionics")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("Wiring")] = df['Mass (kg)'][row + s_items.index("Wiring")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("Pressurant Tank")] = df['Mass (kg)'][row + s_items.index("Pressurant Tank")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("SRM Dome Top")] = df['Mass (kg)'][row + s_items.index("SRM Dome Top")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("Solid Propellant Casing")] = df['Mass (kg)'][row + s_items.index("Solid Propellant Casing")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("SRM Dome Bottom")] = df['Mass (kg)'][row + s_items.index("SRM Dome Bottom")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("Solid Propellant Residual")] = 0   
                df['Jroll'][row + s_items.index("Aft Skirt")] = df['Mass (kg)'][row + s_items.index("Aft Skirt")] * pow(step.r,2)   
                df['Jroll'][row + s_items.index("Gimballs")] = 0    
                df['Jroll'][row + s_items.index("Nozzle")] = 0    
                df['Jroll'][row + s_items.index("Solid Propellant")] = 0    
                
                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['Jroll'][row + s_items.index("Nose Cone")] = df['Mass (kg)'][row + s_items.index("Nose Cone")] * pow(step.r,2)
                    df['Jroll'][row + s_items.index("Forward Skirt")] = df['Mass (kg)'][row + s_items.index("Forward Skirt")] * pow(step.r,2)
                elif i == len(self.listOfSteps) - 1: # if last step                  
                    df['Jroll'][row + s_items.index("Forward Skirt")] = df['Mass (kg)'][row + s_items.index("Forward Skirt")] * pow(step.r,2)
                    df['Jroll'][pl_items.index("PAF")] = df['Mass (kg)'][pl_items.index("PAF")] * pow(step.r,2)                   
                    df['Jroll'][pl_items.index("Payload")] =  df['Mass (kg)'][pl_items.index("Payload")] * pow(step.r,2)     
                    df['Jroll'][pl_items.index("PLF")] = df['Mass (kg)'][pl_items.index("PLF")] * pow(step.r,2)    
                
                
                row += len(s_items)
    
    def initJRoll(self, df, pl_items, l_items, s_items):
        row = len(pl_items)
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            
            if step.propulsion == "Liquid":
                df['Jroll'][row + l_items.index("Avionics")] = self.JRollThinRing(df['Mass (kg)'][row + l_items.index("Avionics")], step.r)
                df['Jroll'][row + l_items.index("Wiring")] = 0 # NEEDS EDIT
                df['Jroll'][row + l_items.index("Fuel Dome Top")] = self.JRollThinEllipsoid(df['Mass (kg)'][row + l_items.index("Fuel Dome Top")], step.r, step.dome_f[0])
                df['Jroll'][row + l_items.index("Fuel Cylinder")] = self.JRollThinCyl(df['Mass (kg)'][row + l_items.index("Fuel Cylinder")], step.r)
                df['Jroll'][row + l_items.index("Fuel Dome Bottom")] = self.JRollThinEllipsoid(df['Mass (kg)'][row + l_items.index("Fuel Dome Bottom")], step.r, step.dome_f[0])
                SA_F_Tank = step.cyl_f[1] + 2 * step.dome_f[1]
                m_ins_f = df['Mass (kg)'][row + l_items.index("Fuel Insulation")]
                m_ins_f_dome = step.dome_f[1]/SA_F_Tank*m_ins_f
                m_ins_f_cyl = step.cyl_f[1]/SA_F_Tank*m_ins_f
                if m_ins_f == 0:
                    df['Jroll'][row + l_items.index("Fuel Insulation")] = 0
                else:
                    df['Jroll'][row + l_items.index("Fuel Insulation")] = (2*m_ins_f_dome*self.JRollThinEllipsoid(m_ins_f_dome, step.r, step.dome_f[0]) + m_ins_f_cyl*self.JRollThinCyl(m_ins_f_cyl, step.r))/m_ins_f
                df['Jroll'][row + l_items.index("Fuel Residual")] = 0   
                df['Jroll'][row + l_items.index("Intertank")] = self.JRollThinCyl(df['Mass (kg)'][row + l_items.index("Intertank")], step.r)
                df['Jroll'][row + l_items.index("Ox Dome Top")] = self.JRollThinEllipsoid(df['Mass (kg)'][row + l_items.index("Ox Dome Top")], step.r, step.dome_ox[0])
                df['Jroll'][row + l_items.index("Ox Cylinder")] = self.JRollThinCyl(df['Mass (kg)'][row + l_items.index("Ox Cylinder")], step.r)
                df['Jroll'][row + l_items.index("Ox Dome Bottom")] = self.JRollThinEllipsoid(df['Mass (kg)'][row + l_items.index("Ox Dome Top")], step.r, step.dome_ox[0])
                SA_Ox_Tank = step.cyl_ox[1] + 2 * step.dome_ox[1]
                m_ins_ox = df['Mass (kg)'][row + l_items.index("Ox Insulation")]
                m_ins_ox_dome = step.dome_ox[1]/SA_Ox_Tank*m_ins_ox
                m_ins_ox_cyl = step.cyl_ox[1]/SA_Ox_Tank*m_ins_ox
                df['Jroll'][row + l_items.index("Ox Insulation")] = (2*m_ins_ox_dome*self.JRollThinEllipsoid(m_ins_ox_dome, step.r, step.dome_ox[0]) + m_ins_ox_cyl*self.JRollThinCyl(m_ins_ox_cyl, step.r))/m_ins_ox
                df['Jroll'][row + l_items.index("Ox Residual")] = 0  
                df['Jroll'][row + l_items.index("Pressurant Tank")] = self.JRollThinHemisphere(df['Mass (kg)'][row + l_items.index("Pressurant Tank")], step.r)
                df['Jroll'][row + l_items.index("Aft Skirt")] = self.JRollThinCyl(df['Mass (kg)'][row + l_items.index("Aft Skirt")], step.r)
                df['Jroll'][row + l_items.index("Thrust Structure")] = 0
                df['Jroll'][row + l_items.index("Gimballs")] = 0 
                df['Jroll'][row + l_items.index("Engines")] = 0 
                df['Jroll'][row + l_items.index("Fuel")] = 0 
                df['Jroll'][row + l_items.index("Oxidizer")] = 0 

                if i < len(self.listOfSteps) - 1: # if not last step
                   # Interstage ( using formual for centroid of trapezoid)
                    df['Jroll'][row + l_items.index("Forward Skirt")] = self.JRollThinTrap(df['Mass (kg)'][row + l_items.index("Forward Skirt")], step.r, self.listOfSteps[i+1].r)
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['Jroll'][row + l_items.index("Forward Skirt")] = self.JRollThinCyl(df['Mass (kg)'][row + l_items.index("Forward Skirt")], step.r)
                    df['Jroll'][pl_items.index("PAF")] = self.JRollThinRing(df['Mass (kg)'][pl_items.index("PAF")], step.r)
                    df['Jroll'][pl_items.index("Payload")] = self.JRollSolidCyl(df['Mass (kg)'][pl_items.index("Payload")], step.r)
                    last_step = self.listOfSteps[len(self.listOfSteps)-1]
                    SA_PLF = sum(last_step.fairing[1])
                    m_PLF = df['Mass (kg)'][pl_items.index("PLF")]
                    m_PLF_cyl = step.fairing[1][0]/SA_PLF*m_PLF
                    m_PLF_cone = step.fairing[1][1]/SA_PLF*m_PLF
                    df['Jroll'][pl_items.index("PLF")] = (m_PLF_cyl*self.JRollThinCyl(m_PLF_cyl, step.r) + m_PLF_cone*self.JRollThinCone(m_PLF_cone, step.r))/m_PLF
                
                
                row += len(l_items)

            elif step.propulsion == 'Solid':
                df['Jroll'][row + s_items.index("Avionics")] = df['Mass (kg)'][row + s_items.index("Avionics")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("Wiring")] = df['Mass (kg)'][row + s_items.index("Wiring")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("Pressurant Tank")] = df['Mass (kg)'][row + s_items.index("Pressurant Tank")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("SRM Dome Top")] = df['Mass (kg)'][row + s_items.index("SRM Dome Top")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("Solid Propellant Casing")] = df['Mass (kg)'][row + s_items.index("Solid Propellant Casing")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("SRM Dome Bottom")] = df['Mass (kg)'][row + s_items.index("SRM Dome Bottom")] * pow(step.r,2)
                df['Jroll'][row + s_items.index("Solid Propellant Residual")] = 0   
                df['Jroll'][row + s_items.index("Aft Skirt")] = df['Mass (kg)'][row + s_items.index("Aft Skirt")] * pow(step.r,2)   
                df['Jroll'][row + s_items.index("Gimballs")] = 0    
                df['Jroll'][row + s_items.index("Nozzle")] = 0    
                df['Jroll'][row + s_items.index("Solid Propellant")] = 0    
                
                if i < len(self.listOfSteps) - 1: # if not last step
                    if step.parallel:
                        df['Jroll'][row + s_items.index("Nose Cone")] = df['Mass (kg)'][row + s_items.index("Nose Cone")] * pow(step.r,2)
                    df['Jroll'][row + s_items.index("Forward Skirt")] = df['Mass (kg)'][row + s_items.index("Forward Skirt")] * pow(step.r,2)
                elif i == len(self.listOfSteps) - 1: # if last step                  
                    df['Jroll'][row + s_items.index("Forward Skirt")] = df['Mass (kg)'][row + s_items.index("Forward Skirt")] * pow(step.r,2)
                    df['Jroll'][pl_items.index("PAF")] = df['Mass (kg)'][pl_items.index("PAF")] * pow(step.r,2)                   
                    df['Jroll'][pl_items.index("Payload")] =  df['Mass (kg)'][pl_items.index("Payload")] * pow(step.r,2)     
                    df['Jroll'][pl_items.index("PLF")] = df['Mass (kg)'][pl_items.index("PLF")] * pow(step.r,2)    
                
                
                row += len(s_items)

    # Re-arrange dataframe to represent the rocket from top to bottom. 
    # This fcn could be deleted if the massmoments fcn is modified to 
    # originially arrange the dataframe accordingly.
    # This function also initializes attributes of the launch vehicle as a result of sizing (stage masses)
    def rearrangeDF(self, pl_items, l_items, s_items): 
        num_pl_= len(pl_items)
        num_l = len(l_items)
        num_s = len(s_items)
        df_temp = pd.DataFrame(columns=['Item', 'Height (m)', 'Mass (kg)', 'Distance (m)', 'Moment (kg*m)', 'Thickness (m)', 'Distance from CM (m)', 'J0 (kg m^2)', 'm*CM^2 (kg m^2)', 'Jpitch/yaw', 'Jroll'])
        df_temp = df_temp.append(self.df.iloc[0:4, :])
        index_end = len(self.df)
        for i in range(len(self.listOfSteps))[::-1]:
            step = self.listOfSteps[i]

            if step.propulsion == "Liquid":
                index_start = index_end-num_l
            elif step.propulsion == "Solid":
                index_start = index_end-num_s

            df_step = self.df.iloc[index_start:index_end, :]
            df_temp = df_temp.append(df_step)
            self.mi_actual.append(df_step['Mass (kg)'].sum())
            index_end = index_start
        self.mi_actual.reverse()
        for i in range(len(self.listOfSteps)):
            empty_mass = self.mi_actual[i] - self.mp_actual[i]
            self.mf_actual.append(empty_mass)
        self.df = df_temp
        del df_temp
        del df_step

    def generateTrajReqs(self):
        df = pd.DataFrame(columns=['Stage #', 'Radius', 'Initial Stage Mass', 'Empty Stage Mass', 'Propellant Mass', 'Thrust', 'Structural Mass Ratio', 
                                'Isp'], index = range(self.num_steps))
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            df['Stage #'][i] = i + 1
            df['Radius'][i] = step.r
            df['Initial Stage Mass'][i] = self.mp_actual[i]+self.mf_actual[i]
            df['Empty Stage Mass'][i] = self.mf_actual[i]
            df['Propellant Mass'][i] = self.mp_actual[i]
            df['Thrust'][i] = step.T_SL * step.multiplier
            df['Structural Mass Ratio'][i] = self.sigmas[i]
            df['Isp'][i] = self.engine_Isps[i]
            pass
        df.to_csv('LVTrajectory\\'+self.name + 'TrajReqs.csv', index=False)

    def addSlide(self):
        prs = Presentation() # Create Presentation

        Vehicle_slide_layout = prs.slide_layouts[5] # define "title only" layout
        vehicle_diagram = prs.slides.add_slide(Vehicle_slide_layout) # add vehicle diagram slide
        shapes = vehicle_diagram.shapes # declare shapes for ease of access to the slide's shapes
        shapes.title.text = self.name + ' Diagram'
        if (self.name == "Minerva-1") | (self.name == "Zephyr-1"):
            scale = 0.6
        elif (self.name == "Minerva-2") | (self.name == "Zephyr-2"):
            scale = 0.5
        elif (self.name == "Latona-1") | (self.name == "Latona-2"):
            scale = 1
        left = 1
        top = 2

        # add to-scale lines
        h_total = self.df['Distance (m)'][0]
        #print(h_total)
        # for i in range(ceil(h_total)):
        #     line = shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(scale*left), Inches(scale*top), Inches(scale*left), Inches(scale*top))
        #     left += 1
        for step in self.listOfSteps:
            top -= step.r
            if step.propulsion == 'Liquid':
                if step.parallel:
                    top += step.r*2
                nozzle = shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(left-scale*step.L_n), Inches(top + scale*step.r/2), Inches(scale*step.L_n), Inches(scale*step.d_e))
                nozzle.rotation = 90
                nozzle_fill = nozzle.fill
                nozzle_fill.solid()
                nozzle_fill.fore_color.rgb = RGBColor(90, 90, 90)
                # Draw Aft Skirt
                aft_skirt = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.aft_skirt[0]), Inches(scale*step.r*2))
                # Draw Thrust Structure
                thrust_struct = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.T_struct), Inches(scale*step.r*2))
                thrust_struct_fill = thrust_struct.fill
                thrust_struct_fill.solid()
                thrust_struct_fill.fore_color.rgb = RGBColor(15, 60, 160)
                
                left+=scale*(step.aft_skirt[0] - step.dome_ox[0]) # add aft skirt and subtract dome ox height to set middle position for ox dome bottom
                # Draw Bottom Ox Dome
                dome_ox_bot = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(scale*step.dome_ox[0]*2), Inches(scale*step.r*2))
                dome_ox_bot_fill = dome_ox_bot.fill
                dome_ox_bot_fill.solid()
                dome_ox_bot_fill.fore_color.rgb = RGBColor(80, 160, 220)
                # Draw Intertank
                left += scale*(step.cyl_ox[0] + step.dome_ox[0])
                intertank = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.intertank[0]), Inches(scale*step.r*2))
                # Draw Top Ox Dome
                left += scale*(-step.dome_ox[0])
                dome_ox_top = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(scale*step.dome_ox[0]*2), Inches(scale*step.r*2))
                dome_ox_top_fill = dome_ox_top.fill
                dome_ox_top_fill.solid()
                dome_ox_top_fill.fore_color.rgb = RGBColor(80, 160, 220)
                # Draw Ox Cylinder
                left += scale*(-step.cyl_ox[0] + step.dome_ox[0])
                ox_cyl = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.cyl_ox[0]), Inches(scale*step.r*2))
                ox_cyl_fill = ox_cyl.fill
                ox_cyl_fill.solid()
                ox_cyl_fill.fore_color.rgb = RGBColor(80, 160, 220)
                # Draw Pressurant Tank
                left += scale*(step.cyl_ox[0]/2 - step.press_tank[0]/2) # adjust to middle of oxidizer cylinder
                top += scale*(step.r - step.press_tank[0]/2)
                press_tank = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(scale*step.press_tank[0]), Inches(scale*step.press_tank[0]))
                press_tank_fill = press_tank.fill
                press_tank_fill.solid()
                press_tank_fill.fore_color.rgb = RGBColor(255, 125, 125)
                left -= scale*(step.cyl_ox[0]/2 - step.press_tank[0]/2) # adjust back to bottom of ox 
                top -= scale*(step.r - step.press_tank[0]/2)
                # Draw Bottom Fuel Dome
                left+=scale*(step.cyl_ox[0] + step.intertank[0] - step.dome_f[0])
                dome_f_bot = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(scale*step.dome_f[0]*2), Inches(scale*step.r*2))
                dome_f_bot_fill = dome_f_bot.fill
                dome_f_bot_fill.solid()
                dome_f_bot_fill.fore_color.rgb = RGBColor(230, 150, 50)
                # Draw Interstage or Fwd Skirt
                left += scale*(step.cyl_f[0] + step.dome_f[0])
                if step.step_num != self.num_steps:
                    interstage = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.interstage[0]), Inches(scale*step.r*2))
                else:
                    fwd_skirt = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.fwd_skirt[0]), Inches(scale*step.r*2))
                # Draw Top Fuel Dome
                left += scale*(-step.dome_f[0])
                dome_f_top = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(scale*step.dome_f[0]*2), Inches(scale*step.r*2))
                dome_f_top_fill = dome_f_top.fill
                dome_f_top_fill.solid()
                dome_f_top_fill.fore_color.rgb = RGBColor(230, 150, 50)
                # Draw Fuel Cylinder
                left += scale*(-step.cyl_f[0] + step.dome_f[0])
                f_cyl = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.cyl_f[0]), Inches(scale*step.r*2))
                f_cyl_fill = f_cyl.fill
                f_cyl_fill.solid()
                f_cyl_fill.fore_color.rgb = RGBColor(230, 150, 50)

                # Add to left parameter depending on step (does it have an interstage or a fwd skirt)
                if step.step_num != self.num_steps:
                    left += scale*(step.cyl_f[0] + step.interstage[0])
                    pass
                else:
                    left += scale*(step.cyl_f[0] + step.fwd_skirt[0])
                    pass
                if step.step_num == self.num_steps:
                    # Draw fairing cylinder and cone
                    fairing_cyl = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.fairing[0][0]), Inches(scale*step.r*2))
                    left += scale*(step.fairing[0][0])
                    fairing_cone = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(scale*step.r*2), Inches(scale*step.fairing[0][1]))
                    fairing_cone.rotation = 90
                    #fairing_cone = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(scale*step.r*2), Inches(scale*step.fairing[0][0]))
                elif step.parallel:
                    fairing_cone = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left - scale*step.fairing[0]/8), Inches(top + scale*step.fairing[0]/10), Inches(scale*step.r*2), Inches(scale*step.fairing[0]))
                    fairing_cone.rotation = 90
                    left = 1
                    top -= step.r*2
                #fore_color.theme_color = MSO_THEME_COLOR.ACCENT_4
            
            elif step.propulsion == 'Solid':
                if step.parallel:
                    top += step.r*3
                # add shapes by length and increment 'left'
                
                # Draw Nozzle
                nozzle = shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(left-scale*step.L_n), Inches(top + scale*step.r/2), Inches(scale*step.L_n), Inches(scale*step.d_e))
                nozzle.rotation = 90
                nozzle_fill = nozzle.fill
                nozzle_fill.solid()
                nozzle_fill.fore_color.rgb = RGBColor(90, 90, 90)

                # Draw Aft Skirt
                aft_skirt = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.aft_skirt[0]), Inches(scale*step.r*2))

                #Drawing Thrust Structure
                thrust_struct = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.T_struct), Inches(scale*step.r*2))
                thrust_struct_fill = thrust_struct.fill
                thrust_struct_fill.solid()
                thrust_struct_fill.fore_color.rgb = RGBColor(15, 60, 160)

                # Draw SRM Dome Bottom
                left+=scale*(step.aft_skirt[0] - step.dome_f[0]) # add aft skirt and subtract SRM dome height to set middle position for SRM dome bottom
                dome_srm_bot = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(scale*step.dome_f[0]*2), Inches(scale*step.r*2))
                dome_srm_bot_fill = dome_srm_bot.fill
                dome_srm_bot_fill.solid()
                dome_srm_bot_fill.fore_color.rgb = RGBColor(100, 50, 20)
                
                left += scale*(step.dome_f[0] + step.srm_casing[0])
                if step.step_num != self.num_steps:
                    # Draw Interstage
                    interstage = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.interstage[0]), Inches(scale*step.r*2))
                else:
                    # Draw Fwd_skirt
                    fwd_skirt = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.fwd_skirt[0]), Inches(scale*step.r*2))

                # Draw SRM Dome Top
                left -= scale*step.dome_f[0]
                dome_srm_top = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(scale*step.dome_f[0]*2), Inches(scale*step.r*2))
                dome_srm_top_fill = dome_srm_bot.fill
                dome_srm_top_fill.solid()
                dome_srm_top_fill.fore_color.rgb = RGBColor(200, 80, 80)

                # Draw SRM Casing
                left+=scale*(-step.srm_casing[0] + step.dome_f[0])
                srm_casing = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(scale*step.srm_casing[0]), Inches(scale*step.r*2))
                srm_casing_fill = srm_casing.fill
                srm_casing_fill.solid()
                srm_casing_fill.fore_color.rgb = RGBColor(200, 80, 80)
                left+=scale*(step.srm_casing[0] + step.r/2) # gap of r/2 between dome top and pressure tank

                top+=scale*(step.press_tank[0]/2) # adjust pressure tank to center of rocket
                press_tank = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(scale*step.press_tank[0]), Inches(scale*step.press_tank[0]))
                press_tank_fill = press_tank.fill
                press_tank_fill.solid()
                press_tank_fill.fore_color.rgb = RGBColor(255, 125, 125)
                if step.step_num != self.num_steps: left+= scale*(-step.r/2 + step.interstage[0]) # remove gap between SRM casing and pressure tank, add fwd_skirt length
                else: left+= scale*(-step.r/2 + step.fwd_skirt[0])
                top-=scale*step.press_tank[0]/2 # remove pressure tank adjustment
               

            if step.parallel:
                nose_cone = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(scale*step.r*2), Inches(scale*step.fairing[0]))
                nose_cone.rotation = 90
                left = 1
                top -= step.r*3

            top += step.r

        prs.save('PPSlides\\'+self.name + 'Diagram.pptx')

    def testSlide(self):

        prs = Presentation() # Create Presentation

        title_layout = prs.slide_layouts[5] # define "title only" layout
        test_slide = prs.slides.add_slide(title_layout) # add test slide
        shapes = test_slide.shapes # declare shapes for ease of access to the slide's shapes
        shapes.title.text = self.name + ' Diagram'
        
        scale = 0.8
        left = 1
        top = 2
        gap = 2

        step = self.listOfSteps[0]

        #aft_skirt = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top*step.step_num), Inches(scale*step.aft_skirt[0]), Inches(scale*step.r*2))
        #aft_skirt2 = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top*step.step_num), Inches((scale + 0.2)*step.aft_skirt[0]), Inches((scale+0.2)*step.r*2))
        press_tank = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top*step.step_num), Inches(scale*step.press_tank[0]*2), Inches(scale*step.r*2))
        scale*=1.5
        press_tank = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top*step.step_num-2), Inches(scale*step.press_tank[0]*2), Inches(scale*step.r*2))
        prs.save('PPSlides\\'+self.name + 'Test.pptx')

    # Roll Mass Moment of Inertia Equations Based on Component
    # Thin Cone
    def JRollThinCone(self, m, R):
        return m/2*pow(R,2)
    # Thin Cylinder
    def JRollThinCyl(self, m, R):
        return m*pow(R,2)
    # Thin Trapezoid
    def JRollThinTrap(self, m, R, r):
        return m/2*(pow(R,2)+pow(r,2))
    # Thin Elliptic-Hemisphere
    def JRollThinEllipsoid(self, m, R, h):
        return 2*m*h/(3*R-h)*(pow(R,2)-3/4*R*h+3/20*pow(h,2))
    # Thin Hemisphere
    def JRollThinHemisphere(self, m, R):
        return 2/3*m*pow(R,2)
    # Solid Elliptic-Hemisphere
    def JRollSolidEllipsoid(self, m, R):
        return 2/5*m*pow(R,2)
    # Solid Cylinder
    def JRollSolidCyl(self, m, R):
        return m/2*pow(R,2)
    # Thin Ring
    def JRollThinRing(self, m, R):
        return m*pow(R,2)

    # Pitch and Yaw Mass Moment of Inertia Equations Based on Component
    # Thin Cone
    def J0ThinCone(self, m, R, h):
        return m*(pow(R,2)/4 + pow(h,2)/18)
    # Thin Cylinder
    def J0ThinCyl(self, m, R, h):
        return m*(pow(R,2)/2 + pow(h,2)/12)
    # Thin Trapezoid
    def J0ThinTrap(self, m, R, r, h):
        return m*((pow(R, 2) + pow(r, 2))/4 + pow(h, 2)/18*(1 + 2*r*R/(pow(r+R, 2))))
    # Thin Elliptic Hemisphere
    def J0ThinEllipsoid(self, m, R, h):
        return m*(pow(R,2)/5 + pow(h,2))
    # Thin Hemisphere
    def J0ThinHemisphere(self, m,R):
        return m*pow(R,2)
    # Thin Ring
    def J0ThinRing(self, m, R):
        return m*pow(R,2)/2
    # Solid Cylinder
    def J0SolidCyl(self, m, R, h):
        return m*(pow(R,2)/4 + pow(h,2)/12)
    # Solid Elliptic Hemisphere
    def J0SolidEllipsoid(self, m, R, h):
        return m*(pow(R,2)/5 + pow(h,2))
    # Uniform Negliglible Thickness Solid Rod
    def J0SolidRod(self, m, l):
        return 1/12*m*pow(l,2)